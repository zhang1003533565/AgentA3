import base64
import importlib.util
import io
import json
import unittest

from app.multi_agents.catalog import get_agent_catalog, normalize_agent_name
from app.multi_agents.leader_agent.agent import leader_agent
from app.multi_agents.ppt_to_docx_agent.agent import ppt_to_docx_agent
from app.rag.document_conversion import convert_ppt_to_docx
from app.rag.document_conversion.ppt_converter import PptConversionError


PPTX_AVAILABLE = importlib.util.find_spec("pptx") is not None
DOCX_AVAILABLE = importlib.util.find_spec("docx") is not None


class PptToDocxAgentTest(unittest.TestCase):
    def test_catalog_and_leader_route_ppt_to_docx_agent(self):
        names = {item["name"] for item in get_agent_catalog()["agents"]}
        self.assertIn("ppt_to_docx_agent", names)
        self.assertEqual("ppt_to_docx_agent", normalize_agent_name("ppt转word"))
        plan = leader_agent._plan_with_rules("帮我把 PPTX 转 DOCX，图片要保留")
        self.assertEqual("ppt_to_docx_agent", plan.target_agent)
        self.assertFalse(plan.need_retrieval)

    def test_agent_process_returns_json_usage_contract(self):
        payload = json.loads(ppt_to_docx_agent.process("转换这份 PPTX", []))
        self.assertEqual("ppt_to_docx_agent", payload["agent"])
        self.assertEqual("ppt_to_docx", payload["task"])
        self.assertIn(".ppt", payload["supportedInput"])
        self.assertIn(".pptx", payload["supportedInput"])
        self.assertEqual(".docx", payload["output"])


@unittest.skipUnless(PPTX_AVAILABLE and DOCX_AVAILABLE, "python-pptx/python-docx not installed")
class PptToDocxConversionTest(unittest.TestCase):
    def test_convert_pptx_to_docx_preserves_text_table_and_image_metadata(self):
        from docx import Document
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "栈与队列"
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(6), Inches(1))
        textbox.text_frame.text = "栈是后进先出，队列是先进先出。"
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.2), Inches(4), Inches(1))
        table_shape.table.cell(0, 0).text = "结构"
        table_shape.table.cell(0, 1).text = "特点"
        table_shape.table.cell(1, 0).text = "栈"
        table_shape.table.cell(1, 1).text = "LIFO"
        slide.shapes.add_picture(io.BytesIO(_tiny_png_bytes()), Inches(1), Inches(3.5), width=Inches(1))

        buffer = io.BytesIO()
        presentation.save(buffer)
        result = convert_ppt_to_docx(buffer.getvalue(), "course.pptx")

        self.assertEqual("docx", result["format"])
        self.assertEqual(1, result["slideCount"])
        self.assertEqual(1, result["imageCount"])

        docx_bytes = base64.b64decode(result["contentBase64"])
        document = Document(io.BytesIO(docx_bytes))
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        self.assertIn("栈与队列", text)
        self.assertIn("栈是后进先出", text)
        table_text = "\n".join(cell.text for table in document.tables for row in table.rows for cell in row.cells)
        self.assertIn("LIFO", table_text)

    def test_convert_legacy_ppt_requires_libreoffice(self):
        with self.assertRaises(PptConversionError) as ctx:
            convert_ppt_to_docx(b"legacy ppt bytes", "course.ppt")
        self.assertIn("LibreOffice", str(ctx.exception))


def _tiny_png_bytes() -> bytes:
    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
    )


if __name__ == "__main__":
    unittest.main()
