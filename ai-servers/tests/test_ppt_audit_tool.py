from pathlib import Path

from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Inches, Pt

from app.ppt_generation.pptx_export_qa import validate_exported_pptx
from scripts.ppt_audit.audit_pptx import audit_pptx
from scripts.ppt_audit.audit_template_fill import audit_template_fills
from scripts.ppt_audit.audit_template_structure import audit_templates


def _make_deck(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    title.text = "Metric …"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    duplicate = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(4), Inches(1))
    duplicate.text = "Metric …"
    presentation.save(path)


def test_audit_writes_reusable_json_and_markdown(tmp_path):
    pptx = tmp_path / "bad.pptx"
    output = tmp_path / "audit"
    _make_deck(pptx)

    report = audit_pptx(pptx, output)

    assert report["passed"] is False
    assert report["schemaVersion"] == 1
    assert report["package"]["slideCount"] == 1
    codes = {item["code"] for item in report["findings"]}
    assert "TEMPLATE_PLACEHOLDER_LEAK" in codes
    assert "TEXT_TRUNCATED_MARKER" in codes
    assert (output / "audit.json").is_file()
    assert (output / "audit.md").is_file()


def test_audit_reports_clean_minimal_deck(tmp_path):
    pptx = tmp_path / "good.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
    box.text = "数据结构与算法核心知识"
    presentation.save(pptx)

    report = audit_pptx(pptx)

    assert report["passed"] is True
    assert report["summary"]["errors"] == 0


def test_export_qa_uses_visible_glyphs_for_right_aligned_text(tmp_path):
    """重叠的文本框外接矩形不等于右对齐字形真的发生碰撞。"""
    pptx = tmp_path / "right-aligned-footer.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(10), Inches(6), Inches(0.65), Inches(0.25))
    first.text = "Page"
    first.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    first.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    second = slide.shapes.add_textbox(Inches(10.55), Inches(6), Inches(0.35), Inches(0.25))
    second.text = "24"
    second.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    second.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    presentation.save(pptx)

    report = validate_exported_pptx(pptx)

    assert report["passed"] is True
    assert not any(item["kind"] == "TEXT_OVERLAP" for item in report["errors"])


def test_export_qa_still_rejects_visible_text_collision(tmp_path):
    pptx = tmp_path / "real-text-overlap.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(2), Inches(0.5))
    first.text = "重要标题"
    first.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    first.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    second = slide.shapes.add_textbox(Inches(2.2), Inches(2), Inches(2), Inches(0.5))
    second.text = "覆盖文字"
    second.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    second.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    presentation.save(pptx)

    report = validate_exported_pptx(pptx)

    assert report["passed"] is False
    assert any(item["kind"] == "TEXT_OVERLAP" for item in report["errors"])


def test_template_structure_audit_covers_all_embedded_templates():
    root = Path(__file__).parents[1] / "app" / "ppt_generation" / "assets" / "templates"

    report = audit_templates(root)

    assert report["schemaVersion"] == 1
    assert report["summary"]["templates"] == 8
    assert report["summary"]["layouts"] == 158
    assert report["summary"]["connectorCandidates"] > 0
    assert report["summary"]["inferredConnectorBindings"] > 0
    assert (
        report["summary"]["inferredConnectorBindings"]
        + report["summary"]["unresolvedConnectorCandidates"]
        == report["summary"]["connectorCandidates"]
    )
    assert report["summary"]["repeatedSlotLayouts"] > 0
    assert report["summary"]["dynamicLayouts"] > 0


def test_template_fill_audit_covers_all_layouts_without_structural_failures():
    root = Path(__file__).parents[1] / "app" / "ppt_generation" / "assets" / "templates"

    report = audit_template_fills(root)

    assert report["summary"]["layouts"] == 158
    assert report["summary"]["parseErrors"] == 0
    assert report["summary"]["failed"] == 0
