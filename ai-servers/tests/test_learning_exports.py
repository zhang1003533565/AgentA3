import builtins
import importlib
import os
from pathlib import Path
import subprocess
import sys
import time
import zipfile

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.rag.document_conversion import generated_exporter
from app.rag.document_conversion.generated_exporter import (
    GeneratedExportError,
    export_python_code_lab,
)
from app.rag.document_conversion.presentation_exporter import export_presentation


def valid_code_lab():
    return {
        "sourceCode": (
            "def slice_middle(values):\n"
            "    return values[1:-1]\n\n"
            "if __name__ == '__main__':\n"
            "    print(slice_middle([1, 2, 3, 4]))\n"
        ),
        "testCode": (
            "from lab import slice_middle\n\n"
            "def test_slice_middle():\n"
            "    assert slice_middle([1, 2, 3, 4]) == [2, 3]\n"
        ),
        "objectives": ["理解左闭右开切片", "能够组合起止索引"],
        "steps": ["阅读示例", "运行源码", "完成自测"],
        "expectedOutput": "[2, 3]",
        "verificationCases": [
            {"input": "[1, 2, 3, 4]", "expected": "[2, 3]"},
        ],
        "evidenceIds": ["python-list-slicing-01"],
    }


def valid_outline():
    return {
        "slides": [
            {
                "title": "切片语法",
                "bullets": ["基本形式是 sequence[start:stop:step]", "stop 索引不包含在结果中"],
                "evidenceIds": ["python-list-slicing-01"],
            },
            {
                "title": "负索引与省略边界",
                "bullets": ["-1 表示最后一个元素", "省略 start 或 stop 表示使用默认边界"],
                "evidenceIds": ["python-list-slicing-02"],
            },
            {
                "title": "代码练习",
                "bullets": ["预测列表 [0, 1, 2, 3][1:-1] 的输出", "修改步长观察结果"],
                "evidenceIds": ["python-list-slicing-01", "python-list-slicing-02"],
            },
        ]
    }


def test_code_lab_export_contains_runnable_source_guide_tests_and_archive(tmp_path, monkeypatch):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))

    result = export_python_code_lab(
        valid_code_lab(),
        {"title": "列表切片实验", "reviewStatus": "passed"},
    )

    assert {item["ext"] for item in result.attachments} == {"py", "md", "zip"}
    source_attachment = next(item for item in result.attachments if item["ext"] == "py")
    source_path = tmp_path / source_attachment["storageKey"]
    compile(source_path.read_text("utf-8"), "lab.py", "exec")
    guide_path = tmp_path / next(
        item["storageKey"] for item in result.attachments if item["ext"] == "md"
    )
    assert "左闭右开切片" in guide_path.read_text("utf-8")

    archive_path = tmp_path / next(
        item["storageKey"] for item in result.attachments if item["ext"] == "zip"
    )
    with zipfile.ZipFile(archive_path) as archive:
        names = archive.namelist()
        assert any(name.endswith(".py") for name in names)
        assert any("test" in name and name.endswith(".py") for name in names)
        test_name = next(name for name in names if "test" in name and name.endswith(".py"))
        compile(archive.read(test_name).decode("utf-8"), "test_lab.py", "exec")


def test_code_lab_export_accepts_the_typed_workflow_payload(tmp_path, monkeypatch):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = {
        "kind": "code_lab",
        "codeLab": {
            "codeBlocks": [
                {"language": "python", "code": "values = [1, 2, 3]\nprint(values[1:])"},
            ],
            "objectives": ["理解切片"],
            "evidenceIds": ["python-list-slicing-01"],
        },
    }

    result = export_python_code_lab(
        payload,
        {"title": "工作流代码实验", "reviewStatus": "passed"},
    )

    source_path = tmp_path / next(
        item["storageKey"] for item in result.attachments if item["ext"] == "py"
    )
    assert "print(values[1:])" in source_path.read_text("utf-8")


@pytest.mark.parametrize(
    "source",
    [
        "import subprocess\nsubprocess.run(['whoami'])\n",
        "import requests\nrequests.get('https://example.com')\n",
        "eval(input())\n",
        "danger = eval\nprint(danger('1 + 1'))\n",
        "import importlib\nimportlib.import_module('os')\n",
        "import typing\nprint(typing.sys.modules)\n",
        "import math\nmodule_alias = math\nprint(module_alias.sqrt(9))\n",
        "license._Printer__filenames = ['/etc/passwd']\nlicense()\n",
        (
            "def f():\n"
            "    pass\n"
            "print('{0.__globals__[__builtins__]}'.format(f))\n"
        ),
        (
            "from random import randint\n"
            "print('{0.__func__.__globals__[_os].environ}'.format(randint))\n"
        ),
        (
            "f = lambda: None\n"
            "loader = f.__globals__['__builtins__']['__import__']\n"
            "module = loader('subprocess')\n"
            "invoke = module.__getattribute__('run')\n"
            "invoke(['whoami'])\n"
        ),
        (
            "u = '_' * 2\n"
            "generator = (value for value in [])\n"
            "runtime_builtins = generator.gi_frame.f_globals[u + 'builtins' + u]\n"
            "loader = runtime_builtins[u + 'import' + u]\n"
            "module = loader('os')\n"
            "dynamic_getattr = runtime_builtins['getattr']\n"
            "invoke = dynamic_getattr(module, 'system')\n"
            "invoke('whoami')\n"
        ),
        (
            "from re import compile\n"
            "del compile\n"
            "print(compile('1 + 1', '<lesson>', 'eval'))\n"
        ),
        (
            "from re import compile\n"
            "try:\n"
            "    raise ValueError('lesson')\n"
            "except ValueError as compile:\n"
            "    pass\n"
            "print(compile('1 + 1', '<lesson>', 'eval'))\n"
        ),
        "from math import sqrt as open\nprint(open)\n",
        "from math import sqrt as __import__\nprint(__import__)\n",
        "from math import sqrt as eval\nprint(eval)\n",
        "open('/tmp/credential', 'w').write('secret')\n",
    ],
)
def test_code_lab_export_rejects_system_network_dynamic_and_credential_operations(
    tmp_path,
    monkeypatch,
    source,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = source

    with pytest.raises(GeneratedExportError):
        export_python_code_lab(payload, {"title": "危险实验", "reviewStatus": "passed"})

    assert list(tmp_path.iterdir()) == []


def test_code_lab_export_allows_explicit_safe_standard_library_capabilities(
    tmp_path,
    monkeypatch,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = (
        "import math\n"
        "from statistics import mean\n"
        "values = [1, 4, 9]\n"
        "print(math.sqrt(values[-1]), mean(values))\n"
    )

    result = export_python_code_lab(
        payload,
        {"title": "安全标准库实验", "reviewStatus": "passed"},
    )

    assert {item["ext"] for item in result.attachments} == {"py", "md", "zip"}


def test_code_lab_export_allows_core_teaching_methods_static_format_and_safe_aliases(
    tmp_path,
    monkeypatch,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = (
        "from functools import partial\n"
        "from math import sqrt\n\n"
        "class Lesson:\n"
        "    def run(self):\n"
        "        return sqrt(16)\n\n"
        "items = [1, 2, 3]\n"
        "items.remove(2)\n"
        "message = 'hello {}'.format('student')\n"
        "root = sqrt\n"
        "emit = partial(print, message)\n"
        "print(items, Lesson().run(), root(9))\n"
        "emit()\n"
    )

    result = export_python_code_lab(
        payload,
        {"title": "Python 核心语法实验", "reviewStatus": "passed"},
    )

    assert {item["ext"] for item in result.attachments} == {"py", "md", "zip"}


@pytest.mark.parametrize(
    "source",
    [
        (
            "def first(value):\n"
            "    template = 'first {}'\n"
            "    return template.format(value)\n\n"
            "def second(value):\n"
            "    template = 'second {}'\n"
            "    return template.format(value)\n\n"
            "print(first('a'), second('b'))\n"
        ),
        (
            "first, second = ('{}', '{name}')\n"
            "print(first.format('a'), second.format(name='b'))\n"
        ),
        (
            "def render(value):\n"
            "    template = 'hello {}'\n"
            "    return template.format(value)\n\n"
            "def consume(template):\n"
            "    return template.upper()\n\n"
            "print(render('student'), consume('lesson'))\n"
        ),
        (
            "class Lesson:\n"
            "    def render(self, value):\n"
            "        template = 'hello {}'\n"
            "        return template.format(value)\n\n"
            "print(Lesson().render('student'))\n"
        ),
    ],
    ids=[
        "F01-independent-functions",
        "F02-tuple-unpack",
        "F03-sibling-parameter",
        "F04-method-local-template",
    ],
)
def test_code_lab_export_allows_scope_correct_static_format_templates(
    tmp_path,
    monkeypatch,
    source,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = source

    result = export_python_code_lab(
        payload,
        {"title": "字符串格式化实验", "reviewStatus": "passed"},
    )

    assert {item["ext"] for item in result.attachments} == {"py", "md", "zip"}


@pytest.mark.parametrize(
    "source",
    [
        (
            "template = '{}'\n"
            "def rewrite():\n"
            "    global template\n"
            "    template = '{0.value}'\n"
            "rewrite()\n"
            "print(template.format(object()))\n"
        ),
        (
            "def outer():\n"
            "    template = '{}'\n"
            "    def rewrite():\n"
            "        nonlocal template\n"
            "        template = '{0.value}'\n"
            "    rewrite()\n"
            "    return template.format(object())\n"
            "print(outer())\n"
        ),
    ],
    ids=["SEC16A-global-rewrite", "SEC16B-nonlocal-rewrite"],
)
def test_code_lab_export_rejects_cross_scope_format_template_rewrites(
    tmp_path,
    monkeypatch,
    source,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = source

    with pytest.raises(GeneratedExportError):
        export_python_code_lab(
            payload,
            {"title": "跨作用域格式化实验", "reviewStatus": "passed"},
        )

    assert list(tmp_path.iterdir()) == []


@pytest.mark.parametrize(
    "source",
    [
        (
            "template = input()\n"
            "class Lesson:\n"
            "    template = '{}'\n"
            "    values = [template.format(value) for value in [1]]\n"
        ),
        (
            "template = '{}'\n"
            "values = [(template := input()) for value in [1]]\n"
            "print(template.format('value'))\n"
        ),
        (
            "template = '{}'\n"
            "def render(values=[(template := input()) for value in [1]]):\n"
            "    return values\n"
            "print(template.format('value'))\n"
        ),
        (
            "template = '{}'\n"
            "render = lambda values=[(template := input()) for value in [1]]: values\n"
            "print(template.format('value'))\n"
        ),
        (
            "template = input()\n"
            "def render(value=template.format('value')):\n"
            "    template = '{}'\n"
            "    return value\n"
        ),
        (
            "template = input()\n"
            "def decorate(value):\n"
            "    return lambda function: function\n"
            "@decorate(template.format('value'))\n"
            "def render():\n"
            "    template = '{}'\n"
        ),
        (
            "template = input()\n"
            "def render(value: template.format('value')):\n"
            "    template = '{}'\n"
            "    return value\n"
        ),
        (
            "template = input()\n"
            "class Lesson(template.format('value')):\n"
            "    template = '{}'\n"
        ),
        (
            "template = input()\n"
            "def decorate(value):\n"
            "    return lambda lesson: lesson\n"
            "@decorate(template.format('value'))\n"
            "class Lesson:\n"
            "    template = '{}'\n"
        ),
        (
            "template = input()\n"
            "class Lesson(object, label=template.format('value')):\n"
            "    template = '{}'\n"
        ),
        (
            "template = input()\n"
            "class Lesson:\n"
            "    message = template.format('value')\n"
            "    template = '{}'\n"
        ),
    ],
    ids=[
        "SEC17-class-comprehension",
        "SEC18-comprehension-named-expression",
        "SEC18-function-header-comprehension",
        "SEC18-lambda-header-comprehension",
        "SEC19-function-default",
        "SEC19-function-decorator",
        "SEC19-function-annotation",
        "SEC20-class-base",
        "SEC20-class-decorator",
        "SEC20-class-keyword",
        "SEC21-class-temporal-lookup",
    ],
)
def test_code_lab_export_rejects_templates_from_non_runtime_lexical_scopes(
    tmp_path,
    monkeypatch,
    source,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = source

    with pytest.raises(GeneratedExportError):
        export_python_code_lab(
            payload,
            {"title": "词法作用域实验", "reviewStatus": "passed"},
        )

    assert list(tmp_path.iterdir()) == []


def test_code_lab_export_allows_normal_python_course_language_features(
    tmp_path,
    monkeypatch,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = (
        "import re\n"
        "from functools import lru_cache, wraps\n\n"
        "counter = 0\n"
        "def bump():\n"
        "    global counter\n"
        "    counter += 1\n\n"
        "def lesson(value):\n"
        "    return helper(value)\n\n"
        "def helper(value):\n"
        "    return value + 1\n\n"
        "def even(value):\n"
        "    return value == 0 or odd(value - 1)\n\n"
        "def odd(value):\n"
        "    return value > 0 and even(value - 1)\n\n"
        "def outer():\n"
        "    def first(value):\n"
        "        return second(value)\n"
        "    def second(value):\n"
        "        return value + 1\n"
        "    return first(1)\n\n"
        "def tagged(label):\n"
        "    def decorate(function):\n"
        "        @wraps(function)\n"
        "        def wrapped(*args, **kwargs):\n"
        "            return label, function(*args, **kwargs)\n"
        "        return wrapped\n"
        "    return decorate\n\n"
        "cache8 = lru_cache(maxsize=8)\n"
        "@cache8\n"
        "@tagged('python')\n"
        "def cached(value):\n"
        "    return value * 2\n\n"
        "class Client:\n"
        "    \"\"\"讲解 __init__ 与普通实例方法。\"\"\"\n"
        "    def connect(self):\n"
        "        return 'local-demo'\n\n"
        "def apply(function, value):\n"
        "    return function(value)\n\n"
        "template = 'hello {}'\n"
        "pattern = re.compile(r'[a-z]+')\n"
        "print(template.format('student'), apply(helper, 1), Client().connect())\n"
        "print(pattern, outer(), cached(2), even(4))\n"
    )

    result = export_python_code_lab(
        payload,
        {"title": "Python 语言特性实验", "reviewStatus": "passed"},
    )

    assert {item["ext"] for item in result.attachments} == {"py", "md", "zip"}


def test_code_lab_export_allows_trusted_implicit_callbacks_and_decorators(
    tmp_path,
    monkeypatch,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = (
        "from functools import lru_cache, partial, reduce\n"
        "from itertools import starmap\n"
        "from re import sub\n\n"
        "@lru_cache(maxsize=8)\n"
        "def square(value):\n"
        "    return value * value\n\n"
        "@lru_cache\n"
        "def cube(value):\n"
        "    return value * value * value\n\n"
        "@lru_cache(4)\n"
        "def fourth_power(value):\n"
        "    return value * value * value * value\n\n"
        "pairs = [(1, 2), (3, 4)]\n"
        "emit = partial(print, 'value')\n"
        "print(list(map(square, [1, 2])))\n"
        "print(list(map(cube, [1, 2])))\n"
        "print(list(map(fourth_power, [1, 2])))\n"
        "print(sorted([2, 1], key=square))\n"
        "print(reduce(lambda left, right: left + right, [1, 2]))\n"
        "print(list(starmap(lambda left, right: left + right, pairs)))\n"
        "print(sub('a', lambda match: 'b', 'a'))\n"
        "emit()\n"
    )

    result = export_python_code_lab(
        payload,
        {"title": "可信回调实验", "reviewStatus": "passed"},
    )

    assert {item["ext"] for item in result.attachments} == {"py", "md", "zip"}


def test_code_lab_export_allows_trusted_deferred_callback_capabilities(
    tmp_path,
    monkeypatch,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = (
        "from collections import defaultdict\n"
        "from dataclasses import field\n"
        "from functools import cached_property\n"
        "from heapq import nlargest\n"
        "from itertools import accumulate\n\n"
        "class Lesson:\n"
        "    @cached_property\n"
        "    def title(self):\n"
        "        return 'Python'\n\n"
        "values = defaultdict(lambda: 0)\n"
        "scores = list(accumulate([1, 2], lambda left, right: left + right))\n"
        "best = nlargest(1, scores, key=abs)\n"
        "descriptor = property(lambda self: 1)\n"
        "factory = field(default_factory=list)\n"
        "sentinel_iterator = iter(lambda: 1, 1)\n"
        "print(values, best, descriptor, factory, sentinel_iterator)\n"
    )

    result = export_python_code_lab(
        payload,
        {"title": "可信延迟回调实验", "reviewStatus": "passed"},
    )

    assert {item["ext"] for item in result.attachments} == {"py", "md", "zip"}


@pytest.mark.parametrize(
    "source",
    [
        "generator = (value for value in [])\nprint(generator.gi_frame)\n",
        "def inspect_frame(frame):\n    return frame.f_globals\n",
        "def inspect_traceback(traceback):\n    return traceback.tb_frame\n",
    ],
)
def test_code_lab_export_rejects_each_runtime_reflection_boundary_independently(
    tmp_path,
    monkeypatch,
    source,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = source

    with pytest.raises(GeneratedExportError):
        export_python_code_lab(
            payload,
            {"title": "运行时边界实验", "reviewStatus": "passed"},
        )

    assert list(tmp_path.iterdir()) == []


@pytest.mark.parametrize(
    "source",
    [
        "from itertools import pairwise\nprint(list(pairwise([1, 2, 3])))\n",
        "print(aiter([]))\n",
        "print(anext(aiter([])))\n",
    ],
)
def test_code_lab_export_rejects_capabilities_missing_from_python_39(
    tmp_path,
    monkeypatch,
    source,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = source

    with pytest.raises(GeneratedExportError):
        export_python_code_lab(
            payload,
            {"title": "Python 3.9 兼容实验", "reviewStatus": "passed"},
        )

    assert list(tmp_path.iterdir()) == []


def test_code_lab_capability_allowlist_exists_on_the_python_39_runtime():
    missing = {
        module_name: sorted(
            name
            for name in names
            if not hasattr(importlib.import_module(module_name), name)
        )
        for module_name, names in generated_exporter._ALLOWED_PYTHON_MODULE_APIS.items()
    }
    missing = {module_name: names for module_name, names in missing.items() if names}
    missing_builtins = sorted(
        name
        for name in generated_exporter._ALLOWED_PYTHON_BUILTINS
        if name != "__name__" and not hasattr(builtins, name)
    )

    assert missing == {}
    assert missing_builtins == []
    assert "pairwise" not in generated_exporter._ALLOWED_PYTHON_MODULE_APIS["itertools"]
    assert {"aiter", "anext"}.isdisjoint(generated_exporter._ALLOWED_PYTHON_BUILTINS)


def test_code_lab_export_removes_partial_files_when_capacity_validation_fails(
    tmp_path,
    monkeypatch,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    monkeypatch.setattr(generated_exporter, "EXPORT_MAX_BYTES", 16)

    with pytest.raises(RuntimeError, match="AI_EXPORT_MAX_BYTES"):
        export_python_code_lab(
            valid_code_lab(),
            {"title": "容量失败实验", "reviewStatus": "passed"},
        )

    assert list(tmp_path.iterdir()) == []


def test_code_lab_export_requires_reviewed_valid_python(tmp_path, monkeypatch):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))

    with pytest.raises(GeneratedExportError, match="审核"):
        export_python_code_lab(valid_code_lab(), {"title": "未审核", "reviewStatus": "pending"})

    payload = valid_code_lab()
    payload["sourceCode"] = "def broken(:\n"
    with pytest.raises(GeneratedExportError, match="语法"):
        export_python_code_lab(payload, {"title": "语法错误", "reviewStatus": "passed"})

    assert list(tmp_path.iterdir()) == []


@pytest.mark.parametrize(
    "source",
    [
        "def lesson(value):\n    global value\n    return value\n",
        "def lesson():\n    nonlocal missing\n    return missing\n",
    ],
    ids=["invalid-global-parameter", "invalid-nonlocal-binding"],
)
def test_code_lab_export_rejects_semantically_uncompilable_python(
    tmp_path,
    monkeypatch,
    source,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    payload = valid_code_lab()
    payload["sourceCode"] = source

    with pytest.raises(GeneratedExportError, match="语法"):
        export_python_code_lab(
            payload,
            {"title": "编译语义实验", "reviewStatus": "passed"},
        )

    assert list(tmp_path.iterdir()) == []


def test_presentation_export_creates_openable_widescreen_pptx_with_evidence_footer(
    tmp_path,
    monkeypatch,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))

    attachment = export_presentation(
        valid_outline(),
        {
            "title": "Python 列表切片",
            "subtitle": "面向初学者的个性化课件",
            "reviewStatus": "passed",
            "evidenceIds": ["python-list-slicing-01", "python-list-slicing-02"],
        },
    )

    assert attachment["ext"] == "pptx"
    presentation = Presentation(tmp_path / attachment["storageKey"])
    assert presentation.slide_width == Inches(13.333)
    assert presentation.slide_height == Inches(7.5)
    assert len(presentation.slides) >= 5
    text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Python 列表切片" in text
    assert "证据：python-list-slicing-01" in text
    assert "参考依据" in text


def test_presentation_export_is_byte_deterministic_for_the_same_input(tmp_path, monkeypatch):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    metadata = {
        "title": "确定性 Python 切片课件",
        "subtitle": "同一输入必须得到同一文件摘要",
        "reviewStatus": "passed",
        "evidenceIds": ["python-list-slicing-01", "python-list-slicing-02"],
    }

    first = export_presentation(valid_outline(), metadata)
    time.sleep(2.1)
    second = export_presentation(valid_outline(), metadata)

    assert first["storageKey"] != second["storageKey"]
    assert first["sha256"] == second["sha256"]
    assert (tmp_path / first["storageKey"]).read_bytes() == (
        tmp_path / second["storageKey"]
    ).read_bytes()


def test_presentation_export_is_byte_deterministic_across_processes(tmp_path):
    project_root = Path(__file__).resolve().parents[1]
    script = """
from app.rag.document_conversion.presentation_exporter import export_presentation

outline = {
    "slides": [
        {
            "title": "切片语法",
            "bullets": ["左闭右开", "支持步长"],
            "evidenceIds": ["python-list-slicing-01"],
        },
        {
            "title": "边界练习",
            "bullets": ["预测输出", "验证负索引"],
            "evidenceIds": ["python-list-slicing-02"],
        },
    ]
}
export_presentation(
    outline,
    {
        "title": "跨进程确定性课件",
        "subtitle": "相同输入保持相同字节",
        "reviewStatus": "passed",
        "evidenceIds": ["python-list-slicing-01", "python-list-slicing-02"],
    },
)
"""
    payloads = []
    for name in ("first", "second"):
        export_root = tmp_path / name
        environment = os.environ.copy()
        environment["AI_EXPORT_ROOT"] = str(export_root)
        environment["PYTHONPATH"] = os.pathsep.join(
            filter(
                None,
                [str(project_root), environment.get("PYTHONPATH", "")],
            )
        )
        subprocess.run(
            [sys.executable, "-c", script],
            cwd=project_root,
            env=environment,
            check=True,
            capture_output=True,
            text=True,
        )
        exported_files = list(export_root.glob("*.pptx"))
        assert len(exported_files) == 1
        payloads.append(exported_files[0].read_bytes())

    assert payloads[0] == payloads[1]


def test_presentation_export_does_not_fetch_remote_images(tmp_path, monkeypatch):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    outline = valid_outline()
    outline["slides"][0]["imageUrl"] = "https://example.invalid/should-not-be-fetched.png"

    def fail_if_network_is_used(*_args, **_kwargs):
        raise AssertionError("presentation export must not fetch remote images")

    monkeypatch.setattr("urllib.request.urlopen", fail_if_network_is_used)

    attachment = export_presentation(
        outline,
        {
            "title": "无远程图片课件",
            "reviewStatus": "passed",
            "evidenceIds": ["python-list-slicing-01"],
        },
    )

    assert (tmp_path / attachment["storageKey"]).is_file()


def test_presentation_export_accepts_legacy_markdown_outline(tmp_path, monkeypatch):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    attachment = export_presentation(
        "## 切片语法\n- 左闭右开\n## 负索引\n- 从序列末尾定位",
        {
            "title": "Markdown 大纲课件",
            "reviewStatus": "passed",
            "evidenceIds": ["python-list-slicing-01"],
        },
    )

    presentation = Presentation(tmp_path / attachment["storageKey"])
    text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "切片语法" in text
    assert "负索引" in text


def test_presentation_export_requires_list_evidence_ids_and_cleans_capacity_failure(
    tmp_path,
    monkeypatch,
):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))
    with pytest.raises(GeneratedExportError, match="证据"):
        export_presentation(
            {"slides": [{"title": "一", "bullets": ["一"]}, {"title": "二", "bullets": ["二"]}]},
            {
                "title": "错误证据类型",
                "reviewStatus": "passed",
                "evidenceIds": "not-a-list",
            },
        )

    monkeypatch.setattr(generated_exporter, "EXPORT_MAX_BYTES", 16)
    with pytest.raises(RuntimeError, match="AI_EXPORT_MAX_BYTES"):
        export_presentation(
            valid_outline(),
            {
                "title": "容量失败课件",
                "reviewStatus": "passed",
                "evidenceIds": ["python-list-slicing-01"],
            },
        )

    assert list(tmp_path.iterdir()) == []


def test_presentation_export_rejects_unreviewed_or_too_short_outline(tmp_path, monkeypatch):
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path))

    with pytest.raises(GeneratedExportError, match="审核"):
        export_presentation(
            valid_outline(),
            {"title": "未审核课件", "reviewStatus": "rejected"},
        )
    with pytest.raises(GeneratedExportError, match="至少 2"):
        export_presentation(
            {"slides": [{"title": "只有一页"}]},
            {"title": "过短课件", "reviewStatus": "passed"},
        )
    with pytest.raises(GeneratedExportError, match="slides.*列表"):
        export_presentation(
            {"slides": "AB"},
            {
                "title": "错误结构课件",
                "reviewStatus": "passed",
                "evidenceIds": ["python-list-slicing-01"],
            },
        )

    assert list(tmp_path.iterdir()) == []
