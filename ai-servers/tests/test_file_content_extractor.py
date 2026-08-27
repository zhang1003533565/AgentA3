import base64
import io
import unittest

from app.rag.document_conversion.file_content_extractor import extract_file_content


PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9Y9ZQmcAAAAASUVORK5CYII="
)


class FileContentExtractorTest(unittest.TestCase):
    def test_markdown_extracts_text_and_embedded_image(self):
        payload = "# 标题\n\n正文\n\n![图](data:image/png;base64,{})".format(base64.b64encode(PNG_1X1).decode("ascii"))
        result = extract_file_content("markdown_to_text_tool", "sample.md", payload.encode("utf-8"))
        self.assertEqual("text_with_images", result["mode"])
        self.assertIn("正文", result["text"])
        self.assertEqual(1, result["imageCount"])

    def test_txt_extracts_plain_text(self):
        result = extract_file_content("txt_to_text_tool", "sample.txt", "第一行\n第二行".encode("utf-8"))
        self.assertEqual("text", result["mode"])
        self.assertEqual(2, result["stats"]["lineCount"])

    def test_word_extracts_text_and_image(self):
        from docx import Document
        from docx.shared import Inches

        document = Document()
        document.add_paragraph("Word 测试正文")
        document.add_picture(io.BytesIO(PNG_1X1), width=Inches(0.2))
        stream = io.BytesIO()
        document.save(stream)
        result = extract_file_content("word_to_text_tool", "sample.docx", stream.getvalue())
        self.assertEqual("text_with_images", result["mode"])
        self.assertIn("Word 测试正文", result["text"])
        self.assertEqual(1, result["imageCount"])

    def test_ppt_extracts_text_and_image(self):
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        text_box.text = "PPT 测试正文"
        slide.shapes.add_picture(io.BytesIO(PNG_1X1), Inches(1), Inches(2), width=Inches(0.2))
        stream = io.BytesIO()
        presentation.save(stream)
        result = extract_file_content("ppt_to_text_tool", "sample.pptx", stream.getvalue())
        self.assertEqual("text_with_images", result["mode"])
        self.assertIn("PPT 测试正文", result["text"])
        self.assertEqual(1, result["imageCount"])

    def test_pdf_extracts_text(self):
        import fitz

        document = fitz.open()
        page = document.new_page()
        page.insert_text((72, 72), "PDF test content")
        payload = document.tobytes()
        document.close()
        result = extract_file_content("pdf_to_text_tool", "sample.pdf", payload)
        self.assertEqual("text", result["mode"])
        self.assertIn("PDF test content", result["text"])

    def test_scanned_pdf_returns_page_image(self):
        import fitz

        document = fitz.open()
        document.new_page(width=100, height=100)
        payload = document.tobytes()
        document.close()
        result = extract_file_content("pdf_to_text_tool", "scan.pdf", payload)
        self.assertEqual("image_only", result["mode"])
        self.assertGreaterEqual(result["imageCount"], 1)


if __name__ == "__main__":
    unittest.main()
