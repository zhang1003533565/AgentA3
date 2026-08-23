from pptx import Presentation
from pptx.util import Inches

from app.ppt_generation.pptx_export_qa import validate_exported_pptx


def test_export_qa_rejects_template_placeholder_and_text_overlap(tmp_path):
    path = tmp_path / "bad.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    first.text = "Our Team"
    second = slide.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(4), Inches(1))
    second.text = "绪论：数据结构基本概念"
    presentation.save(path)

    report = validate_exported_pptx(path)

    assert report["passed"] is False
    assert any(error["kind"] == "TEMPLATE_PLACEHOLDER" for error in report["errors"])
    assert any(error["kind"] == "TEXT_OVERLAP" for error in report["errors"])


def test_export_qa_allows_clean_text_boxes(tmp_path):
    path = tmp_path / "good.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(0.5))
    title.text = "数据结构与算法核心知识详解"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(1))
    body.text = "数据结构研究数据元素之间的逻辑关系、存储方式和基本运算。"
    presentation.save(path)

    report = validate_exported_pptx(path)

    assert report["passed"] is True
    assert report["errors"] == []


def test_export_qa_detects_template_markers_and_footer_page_mismatch(tmp_path):
    path = tmp_path / "template-leak.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    for index, footer in enumerate(("09", "10", "03"), start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        if index == 1:
            marker = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            marker.text = "Metric / Last Year / Revenue"
        page_number = slide.shapes.add_textbox(Inches(12.75), Inches(7.1), Inches(0.25), Inches(0.2))
        page_number.text = footer
    presentation.save(path)

    report = validate_exported_pptx(path)

    assert any(error["kind"] == "TEMPLATE_PLACEHOLDER" for error in report["errors"])
    assert any(warning["kind"] == "PAGE_NUMBER_MISMATCH" for warning in report["warnings"])
