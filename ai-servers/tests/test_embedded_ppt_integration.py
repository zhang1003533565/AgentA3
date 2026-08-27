from __future__ import annotations

import os
from pathlib import Path
from copy import deepcopy

import pytest
from pptx import Presentation

from app.ppt_generation.presenton_html_renderer import render_presenton_html
from app.ppt_generation.service import PptGenerationService
from app.ppt_generation.source_parser import extract_source_text
from app.ppt_generation.template_catalog import EmbeddedTemplateCatalog


def test_embedded_engine_is_default(monkeypatch, tmp_path):
    monkeypatch.setenv("AI_PPT_SOURCE_ROOT", str(tmp_path / "sources"))
    service = PptGenerationService()

    options = service.get_options()

    assert options["engine"] == "presenton-embedded"
    assert options["enhancedEngineAvailable"] is True
    assert options["editorEnabled"] is False
    assert len(options["templates"]) == 8
    assert options["templates"][0]["layouts"]
    assert options["templates"][0]["layouts"][0]["previewTexts"]


def test_template_catalog_reads_bundled_resources():
    catalog = EmbeddedTemplateCatalog()
    ids = {item["id"] for item in catalog.list_templates()}

    assert ids == {"dynamic", "editorial", "executive", "general", "modern", "momentum", "standard", "swift"}
    content, content_type = catalog.thumbnail("general")
    assert content.startswith(b"\x89PNG")
    assert content_type == "image/png"
    with pytest.raises(FileNotFoundError):
        catalog.thumbnail("../general")
    assert len(catalog.layout_summaries("general")) == 12
    assert {"text", "image"}.issubset(set(catalog.layout_summaries("general")[0]["elementTypes"]))
    assert catalog.layout_summaries("general")[0]["previewTexts"]


def test_uploaded_source_file_is_local_and_owner_scoped(monkeypatch, tmp_path):
    monkeypatch.setenv("AI_PPT_SOURCE_ROOT", str(tmp_path / "sources"))
    service = PptGenerationService()

    uploaded = service.upload_source_file("42", "material.txt", "text/plain", "栈与队列".encode())
    owned = service._source_files.get_owned("42", uploaded["fileId"])

    assert uploaded["fileId"].startswith("ppt_file_")
    assert Path(owned["localPath"]).is_file()
    assert "localPath" not in uploaded
    assert service._source_files.get_owned("7", uploaded["fileId"]) is None


def test_source_parser_supports_txt_and_pptx(tmp_path):
    txt = tmp_path / "material.txt"
    txt.write_text("数据结构复习", encoding="utf-8")
    assert extract_source_text(txt) == "数据结构复习"

    pptx = tmp_path / "material.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "栈"
    slide.placeholders[1].text = "后进先出"
    presentation.save(pptx)
    assert "后进先出" in extract_source_text(pptx)


def test_source_parser_preserves_docx_headings_and_tables(tmp_path):
    from docx import Document

    docx = tmp_path / "material.docx"
    document = Document()
    document.add_heading("数据结构", level=1)
    document.add_paragraph("数据结构是数据元素之间关系的集合。")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "算法"
    table.cell(0, 1).text = "平均时间"
    table.cell(1, 0).text = "归并排序"
    table.cell(1, 1).text = "O(n log n)"
    document.save(docx)

    text = extract_source_text(docx)

    assert "# 数据结构" in text
    assert "[表格 1]" in text
    assert "归并排序 | O(n log n)" in text


@pytest.mark.parametrize("template_id", [
    "dynamic", "editorial", "executive", "general", "modern", "momentum", "standard", "swift",
])
def test_presenton_html_renderer_creates_pptx_and_previews_without_pdf(monkeypatch, tmp_path, template_id):
    if os.name == "nt":
        pytest.skip("Windows 仅提供预览；官方 PPTX 导出需要 Linux/WSL/Docker 运行时")
    monkeypatch.setenv("AI_EXPORT_ROOT", str(tmp_path / "exports"))
    monkeypatch.setenv("PRESENTON_ENABLE_PPTX", "true")
    catalog = EmbeddedTemplateCatalog()
    layouts = catalog.load(template_id)["layouts"]
    slides = [
        {"type": "cover", "title": "数据结构复习", "content": ["期末核心知识梳理"]},
        {"type": "content", "title": "栈", "content": ["后进先出", "入栈与出栈", "典型应用"]},
        {"type": "content", "title": "队列", "content": ["先进先出", "顺序队列", "循环队列", "链式队列"]},
    ]
    for index, slide in enumerate(slides):
        layout = layouts[index % len(layouts)]
        slide.update({
            "index": index + 1,
            "templateLayoutId": layout["id"],
            "ui": {"components": deepcopy(layout["components"]), "background": "#FFFFFF"},
        })

    attachment, path, previews, pptx_attachment = render_presenton_html(
        slides,
        "数据结构复习",
        {"templateId": template_id, "pptxOnly": True},
    )

    assert attachment is None
    assert path is None
    assert len(previews) == 3
    assert pptx_attachment["type"] == "pptx"
    assert pptx_attachment["templateId"] == template_id


def test_task_cancel_is_persisted_and_owner_scoped(monkeypatch, tmp_path):
    monkeypatch.setenv("AI_PPT_SOURCE_ROOT", str(tmp_path / "sources"))
    service = PptGenerationService()
    task = {
        "taskId": "ppt_task_0123456789abcdef0123456789abcdef",
        "userId": "42",
        "status": "running",
        "stage": "rendering",
        "progress": 50,
        "attachments": [],
        "previews": [],
    }
    service._task_store.put(task)

    assert service.cancel_task("42", task["taskId"])["status"] == "cancelled"
    with pytest.raises(Exception) as error:
        service.cancel_task("7", task["taskId"])
    assert error.value.status_code == 403
