import base64
import io
import re
from pathlib import Path
from typing import Any, Dict, Iterable, List


class PptConversionError(RuntimeError):
    def __init__(self, message: str, status_code: int = 500):
        super().__init__(message)
        self.status_code = status_code


_PLACEHOLDER_TEXTS = frozenset({
    "单击此处添加文本",
    "点击此处添加文本",
    "单击此处添加标题",
    "点击此处添加标题",
    "click to add text",
    "click to add title",
})

_MONOSPACE_FONTS = frozenset({
    "consolas",
    "courier new",
    "courier",
    "cascadia mono",
    "cascadia code",
    "menlo",
    "monaco",
    "monospace",
    "dejavu sans mono",
    "liberation mono",
    "source code pro",
    "jetbrains mono",
    "fira code",
    "inconsolata",
    "sf mono",
    "ubuntu mono",
    "noto sans mono",
})

_CODE_SYNTAX_RE = re.compile(
    r"^(?:def|class|function|var|let|const|import|from|return|if|else|elif|for|while|print|echo|printf|"
    r"public|private|protected|void|int|float|double|char|bool|string|try|except|switch|case)\b|^#|^//|^>>>"
)


def convert_ppt_to_docx(
    ppt_bytes: bytes,
    original_filename: str,
    convert_mode: str = "reflow",
) -> Dict[str, Any]:
    if not ppt_bytes:
        raise PptConversionError("PPTX 文件不能为空", 400)
    if not (original_filename or "").lower().endswith(".pptx"):
        raise PptConversionError("当前仅支持 .pptx 文件；请先将 .ppt 另存为 .pptx 后再上传", 400)
    if not ppt_bytes[:4] == b"PK\x03\x04":
        raise PptConversionError("上传文件不是有效的 PPTX", 400)

    base_name = _safe_stem(original_filename)
    mode = (convert_mode or "").strip().lower()
    if mode == "image":
        # 高清还原模式：PPT → PDF（soffice）→ 每页图片 DOCX
        pdf_result = convert_ppt_to_pdf(ppt_bytes, original_filename)
        pdf_bytes = base64.b64decode(pdf_result["contentBase64"])
        from app.rag.document_conversion.pdf_converter import _convert_to_docx_image
        return _convert_to_docx_image(pdf_bytes, base_name)

    # reflow / editable / 未知值统一走可编辑逻辑（智能编辑兜底）
    try:
        from docx import Document
        from docx.shared import Inches
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception as exc:
        raise PptConversionError("PPTX 转 DOCX 依赖未安装，请安装 python-pptx 和 python-docx", 500) from exc

    try:
        presentation = Presentation(io.BytesIO(ppt_bytes))
    except Exception as exc:
        raise PptConversionError(f"PPTX 解析失败：{exc}", 400) from exc

    document = Document()
    document.add_heading(base_name, level=0)
    image_assets: List[Dict[str, Any]] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = _slide_title(slide)
        if title:
            document.add_heading(title, level=1)
        shapes = sorted(_iter_shapes(slide.shapes), key=lambda shape: (int(getattr(shape, "top", 0) or 0), int(getattr(shape, "left", 0) or 0)))
        wrote_content = False

        for shape in shapes:
            if _is_title_shape(shape, title):
                continue
            if _has_table(shape):
                if _append_table(document, shape):
                    wrote_content = True
                continue
            if _shape_type(shape) == MSO_SHAPE_TYPE.PICTURE:
                asset = _append_picture(document, shape, slide_index, len(image_assets) + 1, Inches)
                if asset:
                    image_assets.append(asset)
                    wrote_content = True
                continue
            if _has_text_frame(shape):
                wrote_content = _append_text_frame(
                    document,
                    shape,
                    is_first_slide=(slide_index == 1),
                    slide_height_emu=int(getattr(presentation, "slide_height", 0) or 0),
                ) or wrote_content

        if not title and not wrote_content:
            document.add_paragraph("本页未检测到可提取的文字、表格或图片。")
        if slide_index < len(presentation.slides):
            document.add_page_break()

    buffer = io.BytesIO()
    document.save(buffer)
    output_bytes = buffer.getvalue()
    return {
        "format": "docx",
        "outputType": "file",
        "downloadType": "file",
        "fileName": f"{base_name}.docx",
        "mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "contentBase64": base64.b64encode(output_bytes).decode("ascii"),
        "contentLength": len(output_bytes),
        "assets": image_assets,
        "imageCount": len(image_assets),
        "slideCount": len(presentation.slides),
        "conversionMode": "pptx_to_docx_reflow",
    }


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _slide_title(slide: Any) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        title = _shape_text(title_shape)
        if title:
            return title
    for shape in slide.shapes:
        text = _shape_text(shape)
        if text:
            return text.splitlines()[0].strip()
    return ""


def _is_title_shape(shape: Any, title: str) -> bool:
    if not title:
        return False
    return _normalize_text(_shape_text(shape)) == _normalize_text(title)


def _shape_type(shape: Any) -> Any:
    try:
        return getattr(shape, "shape_type", None)
    except NotImplementedError:
        return None


def _has_table(shape: Any) -> bool:
    try:
        return bool(getattr(shape, "has_table", False))
    except NotImplementedError:
        return False


def _has_text_frame(shape: Any) -> bool:
    try:
        return bool(getattr(shape, "has_text_frame", False))
    except NotImplementedError:
        return False


def _is_monospace_font(name: str) -> bool:
    """保守判断等宽字体，用于保护代码文本不被垃圾过滤误删。"""
    value = (name or "").strip().lower()
    if not value:
        return False
    if value in _MONOSPACE_FONTS:
        return True
    return "mono" in value or "courier" in value


def _is_code_paragraph(paragraph: Any, text: str) -> bool:
    """判断段落是否应视为代码：等宽字体，或明显缩进 + 代码关键字开头。"""
    runs = list(getattr(paragraph, "runs", None) or [])
    if any(_is_monospace_font(getattr(run.font, "name", None)) for run in runs):
        return True
    value = str(text or "")
    stripped = value.lstrip(" \t")
    if not stripped:
        return False
    if stripped.startswith(">>>"):
        return True
    leading = value[: len(value) - len(stripped)]
    if len(leading) >= 2 and _CODE_SYNTAX_RE.search(stripped):
        return True
    return False


def _heading_level(text: str) -> int:
    """保守识别章节标题层级：6.1→Heading 1、6.1.1→Heading 2、第X章/数字标题→Heading 1；普通正文返回 0。"""
    value = str(text or "").strip()
    if not value or len(value) > 50:
        return 0
    # 三级编号：6.1.1 文件分类 → Heading 2
    if re.match(r"^\d{1,3}\.\d{1,3}\.\d{1,3}\s+\S", value):
        return 2
    # 二级编号：6.1 文件概述 → Heading 1
    if re.match(r"^\d{1,3}\.\d{1,3}\s+\S", value):
        return 1
    # 一级：第X章 xxx / 第 X 章 xxx / 6 文件系统
    if re.match(r"^第\s*[\d一二三四五六七八九十百千万]+\s*章\s+\S", value):
        return 1
    if re.match(r"^\d{1,3}\s+\S", value):
        return 1
    return 0


_MARKDOWN_HEADING_RE = re.compile(r"^(?P<hashes>#{1,3})[ \t]+")
_QUOTE_PREFIX_RE = re.compile(r"^>{1,3}[ \t]+")
_PYTHON_REPL_RE = re.compile(r"^>>>[ \t]+")
_RESIDUE_START_RE = re.compile(r"^[A-Za-z]\s+\d{1,3}\s+[A-Za-z]$")


def _markdown_heading_prefix_len(text: str) -> int:
    """返回行首 Markdown 标题标记长度（# / ## / ### + 空白），无则 0。"""
    match = _MARKDOWN_HEADING_RE.match(str(text or ""))
    return len(match.group(0)) if match else 0


def _markdown_heading_level(text: str) -> int:
    """返回行首 Markdown 标题级别（#→1、##→2、###→3），无则 0。"""
    match = _MARKDOWN_HEADING_RE.match(str(text or ""))
    return len(match.group("hashes")) if match else 0


def _is_residue_start_text(text: str) -> bool:
    """文档开头的已确认残留结构（如 a 23 D）：单字母+数字+单字母，非常保守。"""
    value = str(text or "").strip()
    return bool(_RESIDUE_START_RE.fullmatch(value))


def _is_ppt_page_number(shape: Any, text: str, slide_height_emu: int) -> bool:
    """保守识别 PPT 页码：纯数字 + 页脚占位符/名称/底部位置信号；无法确认则保留。"""
    value = str(text or "").strip()
    if not re.fullmatch(r"\d{1,3}", value):
        return False
    number = int(value)
    if number < 1 or number > 999:
        return False
    try:
        if getattr(shape, "is_placeholder", False):
            ph_type = getattr(getattr(shape, "placeholder_format", None), "type", None)
            if ph_type is not None:
                from pptx.enum.shapes import PP_PLACEHOLDER

                if ph_type in (PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.FOOTER):
                    return True
    except Exception:
        pass
    try:
        name = ((getattr(shape, "name", "") or "") or "").lower()
        if any(keyword in name for keyword in (
            "page number", "pagenum", "slide number", "footer", "页码", "页脚", "页数"
        )):
            return True
    except Exception:
        pass
    if slide_height_emu and slide_height_emu > 0:
        try:
            top = int(getattr(shape, "top", None) or -1)
            height = int(getattr(shape, "height", None) or 0)
            if top >= 0 and height > 0:
                if (top + height) >= slide_height_emu * 0.88 and height <= slide_height_emu * 0.08:
                    return True
        except Exception:
            pass
    return False


def _quote_prefix_len(text: str) -> int:
    """返回行首 Markdown 引用前缀长度（> / >> / >>> + 空白），无则 0。"""
    match = _QUOTE_PREFIX_RE.match(str(text or ""))
    return len(match.group(0)) if match else 0


def _python_repl_prefix_len(text: str) -> int:
    """返回行首 Python REPL 提示符前缀长度（>>> + 空白），无则 0。"""
    match = _PYTHON_REPL_RE.match(str(text or ""))
    return len(match.group(0)) if match else 0


def _strip_markdown_prefix(runs: List[Any], prefix_len: int) -> List[Any]:
    """从 run 序列中删除前 prefix_len 个字符，保留剩余 run 的字体格式。"""
    remaining = prefix_len
    result: List[Any] = []
    for run in runs:
        text = run.text or ""
        if remaining <= 0:
            result.append(run)
            continue
        if remaining >= len(text):
            remaining -= len(text)
            continue
        run.text = text[remaining:]
        remaining = 0
        result.append(run)
    return result


def _chapter_number_level(text: str) -> int:
    """独立章节编号段落：6.1 → Heading 1；6.1.1 → Heading 2；其余返回 0。"""
    value = str(text or "").strip()
    if re.match(r"^\d+\.\d+$", value):
        return 1
    if re.match(r"^\d+\.\d+\.\d+$", value):
        return 2
    return 0


def _should_skip_text(text: str) -> bool:
    """判断段落文本是否为明确垃圾（占位符/Markdown 图片引用/纯符号）。

    内容完整性优先：数字、日期、URL、“第 X 页”等一律保留，不做页码过滤。
    """
    value = str(text or "").strip()
    if not value:
        return True
    normalized = _normalize_text(value)
    if not normalized:
        return True
    lower = normalized.lower()
    if lower in _PLACEHOLDER_TEXTS:
        return True
    # Markdown 图片引用整行：![](xxx) / ![xxx](xxx)
    if re.match(r"^!\[[^\]]*\]\([^)]*\)$", normalized):
        return True
    # 纯符号/标点（无字母、数字、中日韩等文字）
    if not re.search(r"[\w\u4e00-\u9fff\u3040-\u30ff\uac00-\ud7af]", normalized):
        return True
    return False


def _append_text_frame(
    document: Any,
    shape: Any,
    is_first_slide: bool = False,
    slide_height_emu: int = 0,
) -> bool:
    wrote = False
    last_code_paragraph = None
    prev_was_heading = False
    prev_was_chapter_number = False
    for paragraph in shape.text_frame.paragraphs:
        runs = list(paragraph.runs)
        paragraph_text = "".join(run.text for run in runs) if runs else (paragraph.text or "")
        is_code = _is_code_paragraph(paragraph, paragraph_text)
        if not is_code and _should_skip_text(paragraph_text):
            continue
        if is_code:
            repl_len = _python_repl_prefix_len(paragraph_text)
            if repl_len:
                paragraph_text = paragraph_text[repl_len:]
                runs = _strip_markdown_prefix(runs, repl_len)
        else:
            markdown_level = _markdown_heading_level(paragraph_text)
            markdown_len = _markdown_heading_prefix_len(paragraph_text)
            if markdown_len:
                paragraph_text = paragraph_text[markdown_len:]
                runs = _strip_markdown_prefix(runs, markdown_len)
            quote_len = _quote_prefix_len(paragraph_text)
            if quote_len:
                paragraph_text = paragraph_text[quote_len:]
                runs = _strip_markdown_prefix(runs, quote_len)
        chapter_level = _chapter_number_level(paragraph_text) if not is_code else 0
        heading_level = _heading_level(paragraph_text) if not is_code else 0
        if not is_code and markdown_level:
            heading_level = markdown_level
        subtitle_level = 0
        if not is_code and not chapter_level and not heading_level:
            subtitle_level = _subtitle_level(
                shape, paragraph, paragraph_text, after_chapter_number=prev_was_chapter_number
            )
        list_info = None
        if not is_code and not chapter_level and not heading_level and not subtitle_level:
            info = _paragraph_bullet_info(paragraph)
            if info["kind"] != "plain":
                list_info = info

        if is_code:
            doc_paragraph = document.add_paragraph()
            _add_runs(doc_paragraph, runs, paragraph_text)
            _apply_paragraph_format(doc_paragraph, paragraph)
            _apply_code_format(doc_paragraph)
            _set_code_line_spacing(doc_paragraph)
            last_code_paragraph = doc_paragraph
            prev_was_heading = False
            prev_was_chapter_number = False
        else:
            if last_code_paragraph is not None:
                _close_code_block(last_code_paragraph)
                last_code_paragraph = None
            if chapter_level:
                doc_paragraph = document.add_heading(level=chapter_level)
                _add_runs(doc_paragraph, runs, paragraph_text)
                _apply_paragraph_format(doc_paragraph, paragraph)
                _apply_heading_format(doc_paragraph)
                prev_was_heading = True
                prev_was_chapter_number = True
            elif heading_level:
                doc_paragraph = document.add_heading(level=heading_level)
                _add_runs(doc_paragraph, runs, paragraph_text)
                _apply_paragraph_format(doc_paragraph, paragraph)
                _apply_heading_format(doc_paragraph)
                prev_was_heading = True
                prev_was_chapter_number = False
            elif subtitle_level:
                doc_paragraph = document.add_heading(level=subtitle_level)
                _add_runs(doc_paragraph, runs, paragraph_text)
                _apply_paragraph_format(doc_paragraph, paragraph)
                _apply_heading_format(doc_paragraph, space_after=6)
                prev_was_heading = True
                prev_was_chapter_number = False
            elif list_info is not None:
                doc_paragraph = document.add_paragraph()
                _add_runs(doc_paragraph, runs, paragraph_text)
                _apply_paragraph_format(doc_paragraph, paragraph)
                _apply_list_format(document, doc_paragraph, paragraph)
                _apply_list_indent(doc_paragraph, list_info)
                prev_was_heading = False
                prev_was_chapter_number = False
            else:
                if is_first_slide and _is_residue_start_text(paragraph_text):
                    prev_was_heading = False
                    prev_was_chapter_number = False
                    continue
                if _is_ppt_page_number(shape, paragraph_text, slide_height_emu):
                    prev_was_heading = False
                    prev_was_chapter_number = False
                    continue
                quiz_parts = _quiz_parts(paragraph_text)
                if quiz_parts:
                    for part_text in quiz_parts:
                        doc_paragraph = document.add_paragraph()
                        _add_text(doc_paragraph, part_text, runs)
                        _apply_paragraph_format(doc_paragraph, paragraph)
                        _apply_quiz_format(doc_paragraph, part_text)
                        _apply_body_format(doc_paragraph, runs)
                        prev_was_heading = False
                        prev_was_chapter_number = False
                else:
                    doc_paragraph = document.add_paragraph()
                    _add_runs(doc_paragraph, runs, paragraph_text)
                    _apply_paragraph_format(doc_paragraph, paragraph)
                    _apply_quiz_format(doc_paragraph, paragraph_text)
                    _apply_body_format(doc_paragraph, runs)
                    prev_was_heading = False
                    prev_was_chapter_number = False
        wrote = True
    if last_code_paragraph is not None:
        _close_code_block(last_code_paragraph)
    return wrote


def _set_code_line_spacing(doc_paragraph: Any) -> None:
    """代码块内部行：段前段后均为 0，保证连续视觉。"""
    from docx.shared import Pt

    try:
        fmt = doc_paragraph.paragraph_format
        fmt.space_before = Pt(0)
        fmt.space_after = Pt(0)
    except Exception:
        pass


def _close_code_block(doc_paragraph: Any) -> None:
    """代码块最后一行：与后文保留适当间距。"""
    from docx.shared import Pt

    try:
        doc_paragraph.paragraph_format.space_after = Pt(6)
    except Exception:
        pass


def _add_runs(doc_paragraph: Any, runs: List[Any], paragraph_text: str) -> None:
    """按 run 写入段落文本并保留字体格式；无 run 时写入纯文本。"""
    if not runs:
        if paragraph_text:
            doc_paragraph.add_run(paragraph_text)
        return
    for run in runs:
        if not run.text:
            continue
        doc_run = doc_paragraph.add_run(run.text)
        _copy_run_format(doc_run, run)


def _add_text(doc_paragraph: Any, text: str, source_runs: List[Any]) -> None:
    """为拆分段创建文本 run，基准格式取源段落第一个 run。"""
    doc_run = doc_paragraph.add_run(text)
    if source_runs:
        _copy_run_format(doc_run, source_runs[0])


def _apply_heading_format(doc_paragraph: Any, space_after: int = 12) -> None:
    """章节标题视觉格式：粗体 + 段前约 12pt + 段后指定（默认 12pt）。"""
    from docx.shared import Pt

    try:
        for doc_run in doc_paragraph.runs:
            doc_run.font.bold = True
    except Exception:
        pass
    try:
        fmt = doc_paragraph.paragraph_format
        fmt.space_before = Pt(12)
        fmt.space_after = Pt(space_after)
    except Exception:
        pass


def _subtitle_level(shape: Any, paragraph: Any, text: str, after_chapter_number: bool = False) -> int:
    """非常保守地识别子标题 → Heading 2。

    仅当：占位符本身是标题/副标题类型；或 全粗体 + 最大字号≥20pt + 长度≤30；
    或 紧跟纯章节编号之后 + 全粗体 + 最大字号≥18pt + 长度≤30（位置辅助信号）。
    非代码/列表/纯数字。宁可漏识别，不误判正文。
    宁可漏识别，不误判正文。
    """
    value = str(text or "").strip()
    if not value or len(value) > 30:
        return 0
    if re.match(r"^\d{1,3}(\.\d{1,3})*\s*$", value):
        return 0
    try:
        if getattr(shape, "is_placeholder", False):
            ph_type = getattr(getattr(shape, "placeholder_format", None), "type", None)
            if ph_type is not None:
                try:
                    from pptx.enum.shapes import PP_PLACEHOLDER

                    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        return 2
                except Exception:
                    pass
    except Exception:
        pass
    runs = [r for r in (getattr(paragraph, "runs", None) or []) if (r.text or "").strip()]
    if not runs:
        return 0
    bold = all(bool(getattr(r.font, "bold", None)) for r in runs)
    sizes = [r.font.size.pt for r in runs if getattr(r.font, "size", None) is not None]
    if bold and sizes and max(sizes) >= 20:
        return 2
    if after_chapter_number and bold and sizes and max(sizes) >= 18:
        return 2
    return 0


def _apply_body_format(doc_paragraph: Any, runs: List[Any]) -> None:
    """普通正文统一格式：12pt（仅 PPT 未显式设字号时）、无首行/左缩进、段后约 6pt。"""
    from docx.shared import Pt

    try:
        fmt = doc_paragraph.paragraph_format
        fmt.first_line_indent = Pt(0)
        fmt.left_indent = Pt(0)
        fmt.space_after = Pt(6)
    except Exception:
        pass
    try:
        explicit = any(getattr(r.font, "size", None) is not None for r in (runs or []))
        if not explicit:
            for doc_run in doc_paragraph.runs:
                if doc_run.font.size is None:
                    doc_run.font.size = Pt(12)
    except Exception:
        pass


def _apply_code_format(doc_paragraph: Any) -> None:
    """代码段落视觉：强制 Consolas + 浅灰背景 F2F2F2 + 少量左右内边距。"""
    from docx.shared import Cm

    try:
        for doc_run in doc_paragraph.runs:
            doc_run.font.name = "Consolas"
            _set_run_east_asia_font(doc_run, "Consolas")
    except Exception:
        pass
    try:
        fmt = doc_paragraph.paragraph_format
        fmt.left_indent = Cm(0.1)
        fmt.right_indent = Cm(0.1)
    except Exception:
        pass
    _apply_paragraph_shading(doc_paragraph, "F2F2F2")


def _apply_paragraph_shading(doc_paragraph: Any, hex_color: str) -> None:
    """给段落加浅色底纹 w:shd（失败静默跳过）。"""
    try:
        from docx.oxml.ns import qn

        p_pr = doc_paragraph._p.get_or_add_pPr()
        shd = p_pr.makeelement(qn("w:shd"), {})
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:fill"), hex_color)
        p_pr.insert_element_before(
            shd,
            "w:tabs",
            "w:spacing",
            "w:ind",
            "w:jc",
            "w:textAlignment",
            "w:rPr",
            "w:sectPr",
            "w:pPrChange",
        )
    except Exception:
        try:
            p_pr.append(shd)
        except Exception:
            pass


def _apply_list_indent(doc_paragraph: Any, info: Dict[str, Any]) -> None:
    """列表视觉：左缩进约 0.5cm×(层级+1)，悬挂缩进约 0.5cm。"""
    from docx.shared import Cm

    try:
        level = min(max(int(info.get("level", 0) or 0), 0), 3)
        fmt = doc_paragraph.paragraph_format
        fmt.left_indent = Cm(0.5 * (level + 1))
        fmt.first_line_indent = Cm(-0.5)
    except Exception:
        pass


def _strip_hand_marker(doc_paragraph: Any, info: Dict[str, Any]) -> None:
    """去除手写项目符号/编号前缀，避免与 Word 自动列表重复。"""
    if not doc_paragraph.runs:
        return
    first = doc_paragraph.runs[0].text
    if info.get("kind") == "bullet":
        new_text = re.sub(r"^[•●▪▫‣◦·]\s*|^[-–—]\s+", "", first)
    else:
        new_text = re.sub(r"^\d{1,3}[.)、]\s+", "", first)
    if new_text != first:
        doc_paragraph.runs[0].text = new_text


_QUIZ_OPTION_MARKER = re.compile(r"[A-D][.、]")
_QUIZ_PARSE_PREFIXES = ("解析：", "解析:")


def _quiz_parts(text: str) -> List[str]:
    """随堂测试拆分：≥2 个 A./B./C./D. 选项拆分；解析：内容另起一段。

    返回拆分后的段落列表；不满足保守条件返回 None（保持原段落）。
    """
    value = str(text or "").strip()
    if not value:
        return None
    markers = re.findall(r"(?:^|\s)[A-D][.、]", value)
    if len(markers) >= 2:
        parts = [p.strip() for p in re.split(r"(?=\s[A-D][.、])", value) if p.strip()]
        if len(parts) >= 2 and all(re.match(r"^[A-D][.、]", p) for p in parts):
            return parts
        return None
    for prefix in _QUIZ_PARSE_PREFIXES:
        if value.startswith(prefix):
            rest = value[len(prefix):].strip()
            if rest:
                return [prefix, rest]
            return [value]
    return None


def _apply_quiz_format(doc_paragraph: Any, text: str) -> None:
    """随堂测试/解析：加粗。"""
    value = str(text or "").strip()
    if value.startswith("随堂测试"):
        for doc_run in doc_paragraph.runs:
            doc_run.font.bold = True
        return
    for prefix in _QUIZ_PARSE_PREFIXES:
        if value.startswith(prefix):
            _bold_prefix(doc_paragraph, len(prefix))
            return


def _bold_prefix(doc_paragraph: Any, length: int) -> None:
    remaining = length
    for doc_run in doc_paragraph.runs:
        if remaining <= 0:
            break
        doc_run.font.bold = True
        remaining -= len(doc_run.text or "")


def _copy_run_format(doc_run: Any, ppt_run: Any) -> None:
    """尽可能保留 PPT run 的字体信息（字体名/字号/粗斜体/颜色）。"""
    from docx.shared import Pt, RGBColor

    try:
        font = ppt_run.font
    except Exception:
        return
    try:
        name = (font.name or "").strip()
        if name:
            doc_run.font.name = name
            _set_run_east_asia_font(doc_run, name)
    except Exception:
        pass
    try:
        if font.size is not None:
            doc_run.font.size = Pt(font.size.pt)
    except Exception:
        pass
    try:
        if font.bold is not None:
            doc_run.font.bold = bool(font.bold)
    except Exception:
        pass
    try:
        if font.italic is not None:
            doc_run.font.italic = bool(font.italic)
    except Exception:
        pass
    try:
        rgb = font.color.rgb
        if rgb is not None:
            doc_run.font.color.rgb = RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))
    except Exception:
        pass


def _set_run_east_asia_font(doc_run: Any, name: str) -> None:
    """同步设置中文字体（w:eastAsia），避免中文回退字体。"""
    try:
        from docx.oxml.ns import qn

        r_pr = doc_run._element.get_or_add_rPr()
        r_fonts = r_pr.get_or_add_rFonts()
        r_fonts.set(qn("w:eastAsia"), name)
    except Exception:
        pass


def _apply_paragraph_format(doc_paragraph: Any, ppt_paragraph: Any) -> None:
    """保留段落对齐方式。"""
    try:
        alignment = getattr(ppt_paragraph, "alignment", None)
    except Exception:
        return
    if alignment is None:
        return
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from pptx.enum.text import PP_ALIGN

    mapping = {
        PP_ALIGN.LEFT: WD_ALIGN_PARAGRAPH.LEFT,
        PP_ALIGN.CENTER: WD_ALIGN_PARAGRAPH.CENTER,
        PP_ALIGN.RIGHT: WD_ALIGN_PARAGRAPH.RIGHT,
        PP_ALIGN.JUSTIFY: WD_ALIGN_PARAGRAPH.JUSTIFY,
        getattr(PP_ALIGN, "JUSTIFY_LOW", None): WD_ALIGN_PARAGRAPH.JUSTIFY,
        getattr(PP_ALIGN, "DISTRIBUTE", None): WD_ALIGN_PARAGRAPH.JUSTIFY,
    }
    try:
        target = mapping.get(alignment)
        if target is not None:
            doc_paragraph.alignment = target
    except Exception:
        pass


def _paragraph_bullet_info(paragraph: Any) -> Dict[str, Any]:
    """读取 PPT 段落是否显式声明项目符号/编号（buChar/buAutoNum/buNone）。"""
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    info: Dict[str, Any] = {
        "kind": "plain",
        "char": None,
        "level": int(getattr(paragraph, "level", 0) or 0),
    }
    p_pr = paragraph._p.find(ns + "pPr")
    if p_pr is None:
        return info
    if p_pr.find(ns + "buNone") is not None:
        return info
    bu_char = p_pr.find(ns + "buChar")
    if bu_char is not None:
        info["kind"] = "bullet"
        info["char"] = bu_char.get("char") or "•"
        return info
    if p_pr.find(ns + "buAutoNum") is not None:
        info["kind"] = "number"
        return info
    return info


def _apply_list_format(document: Any, doc_paragraph: Any, ppt_paragraph: Any) -> None:
    """将显式 bullet/编号段落映射为 Word 原生列表样式（List Bullet/Number 1-3）。"""
    info = _paragraph_bullet_info(ppt_paragraph)
    if info["kind"] == "plain":
        return
    level = min(max(int(info["level"]), 0), 3)
    if info["kind"] == "bullet":
        style_name = ["List Bullet", "List Bullet 2", "List Bullet 3"][level]
    else:
        style_name = ["List Number", "List Number 2", "List Number 3"][level]
    try:
        doc_paragraph.style = document.styles[style_name]
        _strip_hand_marker(doc_paragraph, info)
        return
    except (KeyError, ValueError):
        pass
    # 兜底：样式缺失时用显式符号 + 缩进保持层级
    indent = 0.22 * min(max(int(info["level"]), 0), 6)
    doc_paragraph.paragraph_format.left_indent = _docx_inches(indent)
    if info["kind"] == "bullet" and doc_paragraph.runs:
        char = info["char"] or "•"
        doc_paragraph.runs[0].text = f"{char} {doc_paragraph.runs[0].text}"


def _append_table(document: Any, shape: Any) -> bool:
    table = getattr(shape, "table", None)
    if table is None or not table.rows or not table.columns:
        return False
    row_count = len(table.rows)
    col_count = len(table.columns)
    if row_count == 0 or col_count == 0:
        return False
    doc_table = document.add_table(rows=row_count, cols=col_count)
    try:
        doc_table.style = "Table Grid"
    except Exception:
        pass
    _apply_table_dimensions(doc_table, table)

    covered = set()
    for row_index in range(row_count):
        for col_index in range(col_count):
            if (row_index, col_index) in covered:
                continue
            cell = table.cell(row_index, col_index)
            span_h = 1
            span_w = 1
            try:
                if cell.is_merge_origin:
                    span_w = int(cell.span_width or 1)
                    span_h = int(cell.span_height or 1)
                elif cell.is_spanned:
                    # 合并续格：无独立内容，跳过
                    covered.add((row_index, col_index))
                    continue
            except Exception:
                span_h = 1
                span_w = 1
            end_r = min(row_index + span_h, row_count)
            end_c = min(col_index + span_w, col_count)
            for rr in range(row_index, end_r):
                for cc in range(col_index, end_c):
                    covered.add((rr, cc))
            doc_cell = doc_table.cell(row_index, col_index)
            if span_h > 1 or span_w > 1:
                try:
                    doc_cell = doc_cell.merge(doc_table.cell(end_r - 1, end_c - 1))
                except Exception:
                    doc_cell = doc_table.cell(row_index, col_index)
            try:
                _write_table_cell(document, doc_cell, cell, is_header=(row_index == 0))
            except Exception:
                # 单个单元格异常不影响整表
                pass
    _apply_table_borders(doc_table)
    return True


def _apply_table_borders(doc_table: Any) -> None:
    """所有单元格四周边框：黑色实线（约 0.75pt）。失败静默跳过。"""
    try:
        from docx.oxml.ns import qn

        for row in doc_table.rows:
            for cell in row.cells:
                tc_pr = cell._tc.get_or_add_tcPr()
                tc_borders = tc_pr.find(qn("w:tcBorders"))
                if tc_borders is None:
                    tc_borders = tc_pr.makeelement(qn("w:tcBorders"), {})
                    tc_pr.insert_element_before(
                        tc_borders,
                        "w:shd",
                        "w:noWrap",
                        "w:tcMar",
                        "w:textDirection",
                        "w:tcFitText",
                        "w:vAlign",
                        "w:hideMark",
                        "w:headers",
                    )
                for edge in ("top", "left", "bottom", "right"):
                    edge_el = tc_borders.find(qn("w:" + edge))
                    if edge_el is None:
                        edge_el = tc_borders.makeelement(qn("w:" + edge), {})
                        tc_borders.append(edge_el)
                    edge_el.set(qn("w:val"), "single")
                    edge_el.set(qn("w:sz"), "6")
                    edge_el.set(qn("w:color"), "000000")
    except Exception:
        pass


def _apply_table_dimensions(doc_table: Any, table: Any) -> None:
    """映射列宽/行高；失败静默跳过，不影响内容。"""
    from docx.shared import Emu

    try:
        for index in range(len(table.columns)):
            width = getattr(table.columns[index], "width", None)
            if width:
                doc_table.columns[index].width = Emu(int(width))
    except Exception:
        pass
    try:
        for index in range(len(table.rows)):
            height = getattr(table.rows[index], "height", None)
            if height:
                doc_table.rows[index].height = Emu(int(height))
    except Exception:
        pass


def _apply_cell_background(doc_cell: Any, ppt_cell: Any) -> None:
    """读取 PPT 单元格纯色背景并写入 Word 底纹；读取失败跳过。"""
    try:
        fill = getattr(ppt_cell, "fill", None)
        if fill is None:
            return
        from pptx.enum.dml import MSO_FILL_TYPE

        fill_type = fill.type
        if fill_type != MSO_FILL_TYPE.SOLID:
            return
        rgb = fill.fore_color.rgb
        if rgb is None:
            return
        hex_color = "%02X%02X%02X" % (int(rgb[0]), int(rgb[1]), int(rgb[2]))
        from docx.oxml.ns import qn

        tc_pr = doc_cell._tc.get_or_add_tcPr()
        shd = tc_pr.makeelement(qn("w:shd"), {})
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:fill"), hex_color)
        try:
            tc_pr.insert_element_before(
                shd,
                "w:noWrap",
                "w:tcMar",
                "w:textDirection",
                "w:tcFitText",
                "w:vAlign",
                "w:hideMark",
                "w:headers",
            )
        except Exception:
            tc_pr.append(shd)
    except Exception:
        pass


def _write_table_cell(document: Any, doc_cell: Any, ppt_cell: Any, is_header: bool = False) -> None:
    """按段写入表格单元格：保留字体/对齐/列表，并做表格安全垃圾过滤；表头行强制加粗。"""
    _apply_cell_background(doc_cell, ppt_cell)
    text_frame = getattr(ppt_cell, "text_frame", None)
    if text_frame is None:
        return
    first_paragraph = True
    for paragraph in text_frame.paragraphs:
        runs = list(paragraph.runs)
        paragraph_text = "".join(run.text for run in runs) if runs else (paragraph.text or "")
        if _should_skip_text(paragraph_text):
            continue
        doc_paragraph = doc_cell.paragraphs[0] if first_paragraph else doc_cell.add_paragraph()
        first_paragraph = False
        if not runs:
            if paragraph_text:
                doc_paragraph.add_run(paragraph_text)
        else:
            for run in runs:
                if not run.text:
                    continue
                doc_run = doc_paragraph.add_run(run.text)
                _copy_run_format(doc_run, run)
        _apply_paragraph_format(doc_paragraph, paragraph)
        _apply_list_format(document, doc_paragraph, paragraph)
        if is_header:
            for doc_run in doc_paragraph.runs:
                doc_run.font.bold = True


def _append_picture(document: Any, shape: Any, slide_index: int, image_index: int, inches_factory: Any) -> Dict[str, Any]:
    """插入 PPT 图片：保留原始比例与尺寸、支持裁剪、异常安全跳过。"""
    try:
        image = getattr(shape, "image", None)
        if image is None:
            return {}
        image_bytes = image.blob
        if not image_bytes:
            return {}
        ext = _safe_extension(getattr(image, "ext", "") or "png")
        name = f"slide-{slide_index}-image-{image_index}.{ext}"
        width_inches, height_inches = _picture_display_size(shape, image)
        usable_width = _page_content_width_inches(document)
        max_width = usable_width * 0.8 if usable_width else None
        if max_width and width_inches > max_width:
            scale = max_width / width_inches
            width_inches = max_width
            height_inches = height_inches * scale
        width_inches = max(width_inches, 0.1)
        height_inches = max(height_inches, 0.1)

        picture = document.add_picture(
            io.BytesIO(image_bytes),
            width=inches_factory(width_inches),
            height=inches_factory(height_inches),
        )
        src_rect = _picture_src_rect(shape)
        if src_rect:
            _apply_docx_picture_crop(picture, src_rect)
        try:
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            document.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
        except Exception:
            pass
        return {
            "name": name,
            "path": f"assets/{name}",
            "type": "image",
            "mimeType": _image_mime_type(ext),
            "slide": slide_index,
            "size": len(image_bytes),
        }
    except Exception:
        # 单张图片异常不影响整篇转换：跳过并继续
        return {}


def _picture_display_size(shape: Any, image: Any) -> tuple:
    """优先使用 PPT shape 原始尺寸（EMU→Inches）；缺失时按图片像素以 96 DPI 估算。"""
    width_emu = float(getattr(shape, "width", 0) or 0)
    height_emu = float(getattr(shape, "height", 0) or 0)
    if width_emu > 0 and height_emu > 0:
        return width_emu / 914400.0, height_emu / 914400.0
    try:
        pixel_w, pixel_h = image.size
    except Exception:
        pixel_w, pixel_h = 0, 0
    if pixel_w > 0 and pixel_h > 0:
        return pixel_w / 96.0, pixel_h / 96.0
    return 4.0, 3.0


def _page_content_width_inches(document: Any) -> float:
    """Word 页面可用宽度（页宽减左右边距），单位 Inches。"""
    try:
        section = document.sections[0]
        usable_emu = section.page_width - section.left_margin - section.right_margin
        if usable_emu and usable_emu > 0:
            return usable_emu / 914400.0
    except Exception:
        pass
    return 6.5


def _picture_src_rect(shape: Any) -> Dict[str, str]:
    """读取 PPT 图片的 a:srcRect（裁剪）原始 l/r/t/b 值，单位与 Word 一致，可直接复制。"""
    try:
        element = getattr(shape, "_element", None)
        if element is None:
            return {}
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        for node in element.iter():
            if node.tag == ns_a + "srcRect":
                attrs: Dict[str, str] = {}
                for key in ("l", "r", "t", "b"):
                    val = node.get(key)
                    if val is not None:
                        try:
                            num = int(float(val))
                        except (TypeError, ValueError):
                            continue
                        if num:
                            attrs[key] = str(num)
                return attrs
    except Exception:
        return {}
    return {}


def _apply_docx_picture_crop(picture: Any, src_rect: Dict[str, str]) -> None:
    """把裁剪信息写入 Word 图片的 a:blipFill/a:srcRect；失败静默降级为原图。"""
    try:
        inline = getattr(picture, "_inline", None)
        if inline is None:
            return
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        blip_fill = None
        for node in inline.iter():
            if str(node.tag).rsplit("}", 1)[-1] == "blipFill":
                blip_fill = node
                break
        if blip_fill is None:
            return
        blip = None
        stretch = None
        for child in list(blip_fill):
            if child.tag == ns_a + "blip":
                blip = child
            elif child.tag == ns_a + "stretch":
                stretch = child
        if blip is None:
            return
        src_rect_el = blip_fill.makeelement(ns_a + "srcRect", src_rect)
        if stretch is not None:
            stretch.addprevious(src_rect_el)
        else:
            blip.addnext(src_rect_el)
    except Exception:
        pass


def _shape_text(shape: Any) -> str:
    if not _has_text_frame(shape):
        return ""
    try:
        return _normalize_text(getattr(shape, "text", "") or "")
    except NotImplementedError:
        return ""


def _normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def _docx_inches(value: float) -> Any:
    from docx.shared import Inches

    return Inches(value)


def _safe_stem(filename: str) -> str:
    stem = Path(filename or "presentation").stem or "presentation"
    safe = re.sub(r'[\\/:*?"<>|\x00-\x1f]+', "-", stem).strip(".- ")
    return safe or "presentation"


def _safe_extension(ext: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9]+", "", ext.lower())
    return safe or "png"


def _image_mime_type(ext: str) -> str:
    return {
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "png": "image/png",
        "gif": "image/gif",
        "webp": "image/webp",
        "bmp": "image/bmp",
        "tif": "image/tiff",
        "tiff": "image/tiff",
    }.get(ext, f"image/{ext}")


def convert_ppt_to_pdf(ppt_bytes: bytes, original_filename: str) -> Dict[str, Any]:
    """PPT/PPTX 转 PDF：使用 LibreOffice(soffice) headless 转换。"""
    import shutil
    import subprocess
    import tempfile

    if not ppt_bytes:
        raise PptConversionError("PPT 文件不能为空", 400)
    ext = Path(original_filename or "presentation.pptx").suffix.lower()
    if ext not in {".ppt", ".pptx"}:
        raise PptConversionError("仅支持 .ppt 或 .pptx 文件", 400)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise PptConversionError("运行环境未安装 LibreOffice，无法生成 PDF", 500)

    base_name = _safe_stem(original_filename)
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir_path = Path(tmp_dir)
        src_path = tmp_dir_path / f"source{ext}"
        src_path.write_bytes(ppt_bytes)
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, str(src_path)],
                check=True,
                capture_output=True,
                timeout=120,
            )
        except subprocess.CalledProcessError as exc:
            detail = exc.stderr.decode("utf-8", errors="ignore") or str(exc)
            raise PptConversionError(f"LibreOffice 转换失败：{detail}", 500) from exc
        except subprocess.TimeoutExpired as exc:
            raise PptConversionError("LibreOffice 转换超时", 500) from exc
        pdf_path = tmp_dir_path / "source.pdf"
        if not pdf_path.is_file():
            raise PptConversionError("LibreOffice 未生成 PDF 文件", 500)
        output_bytes = pdf_path.read_bytes()

    page_count = None
    try:
        import fitz
        pdf_document = fitz.open(stream=output_bytes, filetype="pdf")
        page_count = pdf_document.page_count
        pdf_document.close()
    except Exception:
        page_count = None

    return {
        "format": "pdf",
        "outputType": "file",
        "downloadType": "file",
        "fileName": f"{base_name}.pdf",
        "mimeType": "application/pdf",
        "contentBase64": base64.b64encode(output_bytes).decode("ascii"),
        "contentLength": len(output_bytes),
        "pageCount": page_count,
        "conversionMode": "ppt_to_pdf_libreoffice",
    }
