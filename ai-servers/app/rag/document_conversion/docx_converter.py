import base64
from pathlib import Path
from typing import Any, Dict


class DocxConversionError(RuntimeError):
    def __init__(self, message: str, status_code: int = 500):
        super().__init__(message)
        self.status_code = status_code


def convert_docx_to_pdf(docx_bytes: bytes, original_filename: str) -> Dict[str, Any]:
    """DOCX 转 PDF：使用 LibreOffice(soffice) headless 转换。"""
    import shutil
    import subprocess
    import tempfile

    if not docx_bytes:
        raise DocxConversionError("DOCX 文件不能为空", 400)
    ext = Path(original_filename or "document.docx").suffix.lower()
    if ext != ".docx":
        raise DocxConversionError("仅支持 .docx 文件", 400)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise DocxConversionError("运行环境未安装 LibreOffice，无法生成 PDF", 500)

    base_name = Path(original_filename or "document.docx").stem or "document"
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir_path = Path(tmp_dir)
        src_path = tmp_dir_path / "source.docx"
        src_path.write_bytes(docx_bytes)
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, str(src_path)],
                check=True,
                capture_output=True,
                timeout=120,
            )
        except subprocess.CalledProcessError as exc:
            detail = exc.stderr.decode("utf-8", errors="ignore") or str(exc)
            raise DocxConversionError(f"LibreOffice 转换失败：{detail}", 500) from exc
        except subprocess.TimeoutExpired as exc:
            raise DocxConversionError("LibreOffice 转换超时", 500) from exc
        pdf_path = tmp_dir_path / "source.pdf"
        if not pdf_path.is_file():
            raise DocxConversionError("LibreOffice 未生成 PDF 文件", 500)
        output_bytes = pdf_path.read_bytes()

    page_count = None
    try:
        import fitz
        pdf_document = fitz.open(stream=output_bytes, filetype="pdf")
        page_count = pdf_document.page_count
        pdf_document.close()
    except Exception:
        page_count = None

    return {
        "format": "pdf",
        "outputType": "file",
        "downloadType": "file",
        "fileName": f"{base_name}.pdf",
        "mimeType": "application/pdf",
        "contentBase64": base64.b64encode(output_bytes).decode("ascii"),
        "contentLength": len(output_bytes),
        "pageCount": page_count,
        "conversionMode": "docx_to_pdf_libreoffice",
    }


def convert_docx_to_ppt(
    docx_bytes: bytes,
    original_filename: str,
    convert_mode: str = "image",
) -> Dict[str, Any]:
    """DOCX 转 PPT：LibreOffice 转 PDF 后复用 PDF 转 PPT 图片方案（_convert_to_pptx）。"""
    import shutil
    import subprocess
    import tempfile

    if not docx_bytes:
        raise DocxConversionError("DOCX 文件不能为空", 400)
    ext = Path(original_filename or "document.docx").suffix.lower()
    if ext != ".docx":
        raise DocxConversionError("仅支持 .docx 文件", 400)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise DocxConversionError("运行环境未安装 LibreOffice，无法转换", 500)

    base_name = Path(original_filename or "document.docx").stem or "document"
    if (convert_mode or "").strip().lower() == "smart":
        return _convert_docx_to_smart_pptx(docx_bytes, base_name)

    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir_path = Path(tmp_dir)
        src_path = tmp_dir_path / "source.docx"
        src_path.write_bytes(docx_bytes)
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, str(src_path)],
                check=True,
                capture_output=True,
                timeout=120,
            )
        except subprocess.CalledProcessError as exc:
            detail = exc.stderr.decode("utf-8", errors="ignore") or str(exc)
            raise DocxConversionError(f"LibreOffice 转换失败：{detail}", 500) from exc
        except subprocess.TimeoutExpired as exc:
            raise DocxConversionError("LibreOffice 转换超时", 500) from exc
        pdf_path = tmp_dir_path / "source.pdf"
        if not pdf_path.is_file():
            raise DocxConversionError("LibreOffice 未生成 PDF 文件", 500)
        pdf_bytes = pdf_path.read_bytes()

    # 复用同包 PDF 转 PPT 图片方案，不复制其逻辑
    from app.rag.document_conversion.pdf_converter import _convert_to_pptx

    result = _convert_to_pptx(pdf_bytes, base_name)
    result["conversionMode"] = "docx_to_ppt_page_image"
    return result


def _convert_docx_to_smart_pptx(docx_bytes: bytes, base_name: str) -> Dict[str, Any]:
    """DOCX 转可编辑 PPT：解析标题层级/段落/图片/表格，按结构生成 PPT（Heading1 章节、Heading2 小节、段落要点、图片、表格，超长自动分页）。"""
    import io

    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        raise DocxConversionError("智能生成 PPT 依赖未安装，请安装 python-docx 和 python-pptx", 500) from exc

    try:
        doc = Document(io.BytesIO(docx_bytes))
    except Exception as exc:
        raise DocxConversionError(f"DOCX 解析失败：{exc}", 400) from exc

    image_shapes = []
    for shape in doc.inline_shapes:
        try:
            blip = shape._inline.graphic.graphicData.pic.blipFill.blip
            r_id = blip.get(qn('r:embed'))
            if r_id:
                image_shapes.append(doc.part.related_parts[r_id].blob)
        except Exception:
            continue
    image_cursor = 0
    slides: List[Dict[str, Any]] = []

    def has_content(s: Dict[str, Any]) -> bool:
        return bool(s["title"] or s["subtitle"] or s["bullets"] or s["images"] or s["tables"])

    def start_new_slide() -> Dict[str, Any]:
        slide = {"title": None, "subtitle": None, "bullets": [], "images": [], "tables": []}
        slides.append(slide)
        return slide

    current = start_new_slide()
    for child in doc.element.body.iterchildren():
        if child.tag == qn('w:p'):
            paragraph = Paragraph(child, doc)
            text = (paragraph.text or "").strip()
            style_name = (paragraph.style.name if paragraph.style is not None else "") or ""
            has_image = (
                child.find('.//' + qn('w:drawing')) is not None
                or child.find('.//' + qn('w:pict')) is not None
            )
            if has_image and image_cursor < len(image_shapes):
                if len(current["images"]) >= 1:
                    current = start_new_slide()
                current["images"].append(image_shapes[image_cursor])
                image_cursor += 1
            if text:
                if style_name.startswith('Heading 1') or style_name == 'Title':
                    current = start_new_slide()
                    current["title"] = text
                elif style_name.startswith('Heading 2') or style_name.startswith('Heading 3'):
                    if current["subtitle"] or current["bullets"] or current["images"] or current["tables"]:
                        current = start_new_slide()
                    current["subtitle"] = text
                else:
                    if current["title"] is None and len(slides) == 1 and not has_content(current):
                        current["title"] = text
                    current["bullets"].append(text)
        elif child.tag == qn('w:tbl'):
            table = Table(child, doc)
            if len(current["tables"]) >= 1:
                current = start_new_slide()
            current["tables"].append(table)

    slides = [s for s in slides if has_content(s)]
    if not slides:
        raise DocxConversionError("DOCX 没有可转换的内容", 422)

    # 超长内容自动分页：每页最多 8 条要点
    final_slides: List[Dict[str, Any]] = []
    for s in slides:
        title = s["title"] or (s["subtitle"] if s["subtitle"] else base_name)
        bullets = list(s["bullets"])
        while len(bullets) > 8:
            final_slides.append({"title": title, "subtitle": s["subtitle"], "bullets": bullets[:8], "images": [], "tables": []})
            bullets = bullets[8:]
        final_slides.append({
            "title": title,
            "subtitle": s["subtitle"],
            "bullets": bullets,
            "images": s["images"],
            "tables": s["tables"],
        })

    try:
        presentation = Presentation()
        for s in final_slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content
            title_text = s["title"] or base_name
            if slide.shapes.title is not None:
                slide.shapes.title.text = title_text
            if s["subtitle"]:
                box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
                box.text_frame.text = s["subtitle"]
            content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if content_placeholder is not None and s["bullets"]:
                tf = content_placeholder.text_frame
                tf.text = s["bullets"][0]
                for bullet in s["bullets"][1:]:
                    paragraph = tf.add_paragraph()
                    paragraph.text = bullet
            image_top = 3.2
            for index, image_bytes in enumerate(s["images"]):
                slide.shapes.add_picture(
                    io.BytesIO(image_bytes),
                    Inches(0.5 + index * 4.2),
                    Inches(image_top),
                    height=Inches(3),
                )
            for table in s["tables"]:
                rows = [[cell.text for cell in row.cells] for row in table.rows]
                if rows:
                    cols = max((len(r) for r in rows), default=0)
                    graphic_frame = slide.shapes.add_table(
                        len(rows), cols, Inches(0.5), Inches(3.2), Inches(8), Inches(2.5)
                    )
                    for ri, row in enumerate(rows):
                        for ci, value in enumerate(row):
                            if ci < cols:
                                graphic_frame.table.cell(ri, ci).text = value
    except DocxConversionError:
        raise
    except Exception as exc:
        raise DocxConversionError(f"智能 PPT 生成失败：{exc}", 500) from exc

    buffer = io.BytesIO()
    presentation.save(buffer)
    output_bytes = buffer.getvalue()
    return {
        "format": "pptx",
        "outputType": "file",
        "downloadType": "file",
        "fileName": f"{base_name}.pptx",
        "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "contentBase64": base64.b64encode(output_bytes).decode("ascii"),
        "contentLength": len(output_bytes),
        "pageCount": len(final_slides),
        "imageCount": sum(len(s["images"]) for s in final_slides),
        "conversionMode": "docx_to_ppt_smart",
    }
