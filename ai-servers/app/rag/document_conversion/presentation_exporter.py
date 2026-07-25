import os
import re
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Mapping, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.rag.document_conversion import generated_exporter


_NAVY = RGBColor(15, 35, 64)
_BLUE = RGBColor(44, 107, 237)
_INK = RGBColor(31, 41, 55)
_MUTED = RGBColor(100, 116, 139)
_WHITE = RGBColor(255, 255, 255)


def export_presentation(
    outline: Any,
    metadata: Mapping[str, Any],
) -> Dict[str, Any]:
    """Create a deterministic, local-only PPTX from a reviewed learning outline."""
    generated_exporter._require_passed_review(metadata)
    slide_specs = _normalize_outline(outline)
    if len(slide_specs) < 2:
        raise generated_exporter.GeneratedExportError("课件大纲至少 2 个内容页")

    title = str(metadata.get("title") or "Python 个性化学习课件").strip()
    subtitle = str(metadata.get("subtitle") or "基于课程知识证据生成").strip()
    evidence_ids = _unique_strings(
        _unique_strings(metadata.get("evidenceIds"))
        + [
            evidence_id
            for slide in slide_specs
            for evidence_id in slide.get("evidenceIds", [])
        ]
    )
    if not evidence_ids:
        raise generated_exporter.GeneratedExportError("课件导出缺少课程证据")

    generated_exporter.cleanup_generated_exports(root=generated_exporter._current_export_root())
    path = generated_exporter._new_export_path(generated_exporter._slugify(title), "pptx")
    try:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        presentation.core_properties.created = datetime(2000, 1, 1)
        presentation.core_properties.modified = datetime(2000, 1, 1)
        _add_title_slide(presentation, title, subtitle)
        _add_agenda_slide(presentation, [slide["title"] for slide in slide_specs], evidence_ids)
        for slide in slide_specs:
            _add_content_slide(presentation, slide, evidence_ids)
        _add_references_slide(presentation, evidence_ids)
        generated_exporter._atomic_write_payload(
            path,
            lambda temporary_path: _save_deterministic_presentation(
                presentation,
                temporary_path,
            ),
        )
        attachment = generated_exporter._attachment_for_file(
            path,
            "presentation_export_tool",
            "PowerPoint 课件",
        )
        result = generated_exporter._finalize_export_batch(
            generated_exporter.GeneratedExportResult(
                attachments=[attachment],
                diagnostics={
                    "skipped": False,
                    "contentKind": "presentation",
                    "slideCount": len(presentation.slides),
                    "evidenceIds": evidence_ids,
                    "producedFormats": ["pptx"],
                },
            )
        )
        return result.attachments[0]
    except Exception:
        generated_exporter._delete_export_pair(
            generated_exporter._current_export_root(),
            path.name,
        )
        raise


def _normalize_outline(outline: Any) -> List[Dict[str, Any]]:
    if isinstance(outline, Mapping):
        raw_slides = outline.get("slides") or outline.get("sections") or []
        if not isinstance(raw_slides, list):
            raise generated_exporter.GeneratedExportError(
                "结构化课件的 slides/sections 必须是列表"
            )
    elif isinstance(outline, list):
        raw_slides = outline
    elif isinstance(outline, str):
        return _parse_markdown_outline(outline)
    else:
        raw_slides = []

    slides: List[Dict[str, Any]] = []
    for index, item in enumerate(raw_slides, start=1):
        if isinstance(item, Mapping):
            title = str(item.get("title") or item.get("heading") or f"第 {index} 部分").strip()
            bullets = _normalize_bullets(
                item.get("bullets")
                or item.get("points")
                or item.get("content")
                or item.get("body")
            )
            evidence_ids = _unique_strings(item.get("evidenceIds") or [])
        else:
            title = str(item).strip() or f"第 {index} 部分"
            bullets = []
            evidence_ids = []
        slides.append(
            {
                "title": title[:80],
                "bullets": bullets[:7] or ["结合课程讲解完成本节学习任务。"],
                "evidenceIds": evidence_ids,
            }
        )
    return slides


def _save_deterministic_presentation(
    presentation: Presentation,
    path: Path,
) -> None:
    presentation.save(path)
    _normalize_pptx_zip(path)


def _normalize_pptx_zip(path: Path) -> None:
    with zipfile.ZipFile(path, "r") as source:
        entries = [
            (item.filename, item.compress_type, source.read(item.filename))
            for item in source.infolist()
        ]

    normalized_path = path.with_name(f"{path.name}.normalized")
    try:
        with zipfile.ZipFile(normalized_path, "w") as target:
            for filename, compress_type, content in sorted(entries, key=lambda item: item[0]):
                info = zipfile.ZipInfo(filename, date_time=(1980, 1, 1, 0, 0, 0))
                info.compress_type = compress_type
                info.create_system = 0
                info.external_attr = 0
                info.internal_attr = 0
                info.flag_bits = 0
                target.writestr(
                    info,
                    content,
                    compress_type=compress_type,
                    compresslevel=9,
                )
        os.replace(normalized_path, path)
    finally:
        generated_exporter._safe_unlink(normalized_path)


def _parse_markdown_outline(value: str) -> List[Dict[str, Any]]:
    slides: List[Dict[str, Any]] = []
    current: Optional[Dict[str, Any]] = None
    for raw_line in str(value or "").splitlines():
        line = raw_line.strip()
        heading = re.match(r"^#{2,6}\s+(.+)$", line)
        if heading:
            if current is not None:
                slides.append(current)
            current = {
                "title": heading.group(1).strip()[:80],
                "bullets": [],
                "evidenceIds": [],
            }
            continue
        if current is not None and re.match(r"^[-*]\s+", line):
            current["bullets"].append(re.sub(r"^[-*]\s+", "", line).strip()[:180])
    if current is not None:
        slides.append(current)
    return slides


def _normalize_bullets(value: Any) -> List[str]:
    if isinstance(value, str):
        candidates = [
            re.sub(r"^[-*]\s+", "", line).strip()
            for line in value.splitlines()
            if line.strip()
        ]
    elif isinstance(value, list):
        candidates = [
            str(item.get("text") or item.get("content") or "").strip()
            if isinstance(item, Mapping)
            else str(item).strip()
            for item in value
        ]
    else:
        candidates = []
    return [item[:180] for item in candidates if item]


def _add_title_slide(presentation: Presentation, title: str, subtitle: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _NAVY
    accent = slide.shapes.add_shape(1, Inches(0.75), Inches(1.05), Inches(0.12), Inches(4.7))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _BLUE
    accent.line.fill.background()
    _add_textbox(slide, title, 1.15, 1.3, 11.1, 1.8, 34, _WHITE, bold=True)
    _add_textbox(slide, subtitle, 1.15, 3.35, 10.5, 0.9, 20, RGBColor(203, 213, 225))
    _add_textbox(slide, "A3 · Python 个性化学习资源", 1.15, 5.55, 6.5, 0.5, 12, RGBColor(148, 163, 184))


def _add_agenda_slide(
    presentation: Presentation,
    titles: List[str],
    evidence_ids: List[str],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_light_background(slide)
    _add_slide_title(slide, "学习导航")
    rows = [f"{index:02d}  {title}" for index, title in enumerate(titles, start=1)]
    _add_bullet_box(slide, rows, top=1.55, height=4.9, font_size=22)
    _add_evidence_footer(slide, evidence_ids)


def _add_content_slide(
    presentation: Presentation,
    spec: Mapping[str, Any],
    package_evidence_ids: List[str],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_light_background(slide)
    _add_slide_title(slide, str(spec["title"]))
    _add_bullet_box(slide, list(spec.get("bullets") or []), top=1.55, height=4.95, font_size=24)
    _add_evidence_footer(
        slide,
        list(spec.get("evidenceIds") or []) or package_evidence_ids,
    )


def _add_references_slide(presentation: Presentation, evidence_ids: List[str]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_light_background(slide)
    _add_slide_title(slide, "参考依据")
    rows = [f"课程证据 {index}: {evidence_id}" for index, evidence_id in enumerate(evidence_ids, start=1)]
    _add_bullet_box(slide, rows, top=1.55, height=4.9, font_size=22)
    _add_evidence_footer(slide, evidence_ids)


def _paint_light_background(slide: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _BLUE
    bar.line.fill.background()


def _add_slide_title(slide: Any, title: str) -> None:
    _add_textbox(slide, title, 0.85, 0.55, 11.6, 0.75, 28, _NAVY, bold=True)


def _add_bullet_box(
    slide: Any,
    bullets: List[str],
    *,
    top: float,
    height: float,
    font_size: int,
) -> None:
    shape = slide.shapes.add_textbox(Inches(1.05), Inches(top), Inches(11.15), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(bullet)
        paragraph.level = 0
        paragraph.space_after = Pt(14)
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = _INK
        paragraph.text = f"•  {paragraph.text}"


def _add_evidence_footer(slide: Any, evidence_ids: List[str]) -> None:
    label = "证据：" + "、".join(evidence_ids)
    _add_textbox(slide, label[:240], 0.85, 6.88, 11.55, 0.3, 10, _MUTED, align=PP_ALIGN.RIGHT)


def _add_textbox(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    color: RGBColor,
    *,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return shape


def _unique_strings(values: Any) -> List[str]:
    if not isinstance(values, list):
        return []
    result: List[str] = []
    for value in values:
        normalized = str(value).strip()
        if normalized and normalized not in result:
            result.append(normalized)
    return result


__all__ = ["export_presentation"]
