import base64
import io
from typing import Any, Dict, List

from app.rag.document_conversion.pdf_converter import MAX_PPTX_SLIDES, PdfConversionError

BG_ZOOM = 1.0  # 背景渲染分辨率（约 72 DPI）
BG_JPEG_QUALITY = 85  # 背景 JPEG 压缩质量


def convert_pdf_to_editable_pptx(pdf_bytes: bytes, base_name: str) -> Dict[str, Any]:
    """PDF 转可编辑 PPTX：提取文本块与图片，按坐标生成可编辑元素；空页/扫描页 fallback 为整页图片。"""
    try:
        import fitz
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:
        raise PdfConversionError("PDF 转可编辑 PPTX 依赖未安装，请安装 PyMuPDF 和 python-pptx", 500) from exc

    try:
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception as exc:
        raise PdfConversionError(f"PDF 解析失败：{exc}", 400) from exc

    page_count = pdf_document.page_count
    if page_count <= 0:
        raise PdfConversionError("PDF 没有可转换的页面", 422)
    if page_count > MAX_PPTX_SLIDES:
        raise PdfConversionError(f"PDF 页数超过 {MAX_PPTX_SLIDES} 页，暂不支持转 PPTX", 422)

    try:
        presentation = Presentation()
        slide_layout = presentation.slide_layouts[6]  # blank
        for page_index in range(page_count):
            page = pdf_document.load_page(page_index)
            page_rect = page.rect
            width_inches = max(page_rect.width / 72.0, 0.1)
            height_inches = max(page_rect.height / 72.0, 0.1)
            slide = presentation.slides.add_slide(slide_layout)
            slide.slide_width = Inches(width_inches)
            slide.slide_height = Inches(height_inches)

            # 先提取文字元素（必须在 redaction 之前，redaction 后无法再读取文字）
            text_elements = _extract_text_elements(page)

            # 清理背景：用 redaction 删除文字区域，保留图片与图形
            if text_elements:
                try:
                    for element in text_elements:
                        x0, y0, x1, y1 = element["bbox"]
                        page.add_redact_annot(
                            fitz.Rect(max(x0 - 2, 0), max(y0 - 2, 0), x1 + 2, y1 + 2)
                        )
                    page.apply_redactions(images=2, graphics=1, text=0)
                except Exception:
                    pass  # 清理失败则使用原页背景，不阻断转换

            # 背景层：渲染无文字页面图片铺满幻灯片，保留原始视觉（图片/图形/底纹）
            pixmap = page.get_pixmap(matrix=fitz.Matrix(BG_ZOOM, BG_ZOOM), alpha=False)
            background_bytes = pixmap.tobytes("jpeg", jpg_quality=BG_JPEG_QUALITY)
            slide.shapes.add_picture(
                io.BytesIO(background_bytes),
                0,
                0,
                width=Inches(width_inches),
                height=Inches(height_inches),
            )
            # 编辑层：透明 textbox 覆盖文字区域
            for element in text_elements:
                _add_text_element(slide, element)
    except PdfConversionError:
        raise
    except Exception as exc:
        raise PdfConversionError(f"可编辑 PPTX 生成失败：{exc}", 500) from exc
    finally:
        pdf_document.close()

    buffer = io.BytesIO()
    presentation.save(buffer)
    output_bytes = buffer.getvalue()
    return {
        "format": "pptx",
        "outputType": "file",
        "downloadType": "file",
        "fileName": f"{base_name}.pptx",
        "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "contentBase64": base64.b64encode(output_bytes).decode("ascii"),
        "contentLength": len(output_bytes),
        "pageCount": page_count,
        "imageCount": page_count,
        "conversionMode": "editable",
    }


def _extract_text_elements(page) -> List[Dict[str, Any]]:
    """从页面提取文字块元素（按 PDF 坐标）。"""
    elements: List[Dict[str, Any]] = []
    try:
        data = page.get_text("dict")
    except Exception:
        return elements

    for block in data.get("blocks", []):
        bbox = block.get("bbox")
        if not bbox:
            continue
        if block.get("type") == 0:
            lines = []
            for line in block.get("lines", []):
                spans = []
                for span in line.get("spans", []):
                    text = (span.get("text") or "").strip()
                    if not text:
                        continue
                    spans.append({
                        "text": span.get("text"),
                        "size": float(span.get("size") or 12.0),
                        "font": span.get("font") or "",
                        "bold": bool((span.get("flags") or 0) & 16),
                        "color": span.get("color"),
                    })
                if spans:
                    lines.append({"bbox": tuple(line.get("bbox") or (0, 0, 0, 0)), "spans": spans})
            if lines:
                elements.append({"bbox": tuple(bbox), "lines": lines})
    return elements


def _add_text_element(slide, element: Dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    x0, y0, x1, y1 = element["bbox"]
    left = Inches(max(x0 / 72.0, 0))
    top = Inches(max(y0 / 72.0, 0))
    width = Inches(max((x1 - x0) / 72.0, 0.1))
    height = Inches(max((y1 - y0) / 72.0, 0.1))
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    for index, line_item in enumerate(element["lines"]):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        line_bbox = line_item.get("bbox")
        line_height_pt = (line_bbox[3] - line_bbox[1]) if line_bbox else None
        max_size = max((span.get("size") or 12.0 for span in line_item["spans"]), default=12.0)
        if line_height_pt and 0 < line_height_pt <= max_size * 3:
            try:
                paragraph.line_spacing = Pt(line_height_pt)
            except Exception:
                pass
        for span in line_item["spans"]:
            run = paragraph.add_run()
            run.text = span["text"]
            run.font.size = Pt(span["size"])
            if span.get("bold"):
                run.font.bold = True
            if span.get("font"):
                try:
                    run.font.name = span["font"]
                except Exception:
                    pass
            if span.get("color") is not None:
                try:
                    run.font.color.rgb = RGBColor.from_string(format(int(span["color"]), "06X"))
                except Exception:
                    pass
