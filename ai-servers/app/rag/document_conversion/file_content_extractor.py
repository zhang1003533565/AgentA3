import base64
import io
import mimetypes
import re
from pathlib import Path
from typing import Any, Dict, Iterable, List, Tuple


MAX_FILE_BYTES = 25 * 1024 * 1024
MAX_IMAGES = 20
MAX_IMAGE_BYTES = 5 * 1024 * 1024

TOOL_EXTENSIONS = {
    "markdown_to_text_tool": {".md", ".markdown"},
    "txt_to_text_tool": {".txt"},
    "word_to_text_tool": {".docx"},
    "ppt_to_text_tool": {".pptx"},
    "pdf_to_text_tool": {".pdf"},
}


class FileContentExtractionError(ValueError):
    pass


def extract_file_content(tool_name: str, file_name: str, content: bytes) -> Dict[str, Any]:
    normalized_tool = str(tool_name or "").strip()
    extension = Path(str(file_name or "")).suffix.lower()
    allowed = TOOL_EXTENSIONS.get(normalized_tool)
    if allowed is None:
        raise FileContentExtractionError("不支持的文件提取工具")
    if extension not in allowed:
        expected = "、".join(sorted(allowed))
        raise FileContentExtractionError(f"当前工具仅支持 {expected} 文件")
    if not content:
        raise FileContentExtractionError("上传文件不能为空")
    if len(content) > MAX_FILE_BYTES:
        raise FileContentExtractionError("上传文件不能超过 25MB")

    if normalized_tool == "markdown_to_text_tool":
        text, images, references, stats = _extract_markdown(content)
    elif normalized_tool == "txt_to_text_tool":
        text, images, references, stats = _extract_txt(content)
    elif normalized_tool == "word_to_text_tool":
        text, images, references, stats = _extract_docx(content)
    elif normalized_tool == "ppt_to_text_tool":
        text, images, references, stats = _extract_pptx(content)
    else:
        text, images, references, stats = _extract_pdf(content)

    normalized_text = _normalize_text(text)
    mode = "text_with_images" if normalized_text and images else "text" if normalized_text else "image_only" if images else "empty"
    if mode == "empty":
        raise FileContentExtractionError("文件中没有提取到可用文本或图片")
    return {
        "toolName": normalized_tool,
        "fileName": file_name,
        "inputFormat": extension.lstrip("."),
        "mode": mode,
        "text": normalized_text,
        "textLength": len(normalized_text),
        "imageCount": len(images),
        "images": images,
        "imageReferences": references,
        "stats": stats,
    }


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-16", "gb18030"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise FileContentExtractionError("文本文件编码无法识别，请使用 UTF-8、UTF-16 或 GB18030")


def _extract_markdown(content: bytes) -> Tuple[str, List[Dict[str, Any]], List[str], Dict[str, Any]]:
    text = _decode_text(content)
    images: List[Dict[str, Any]] = []
    references: List[str] = []
    pattern = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")
    for index, match in enumerate(pattern.finditer(text), start=1):
        source = match.group(1).strip().strip("<>\"")
        data_match = re.match(r"data:(image/[^;,]+);base64,(.+)", source, re.IGNORECASE | re.DOTALL)
        if data_match:
            try:
                blob = base64.b64decode(data_match.group(2), validate=True)
                _append_image(images, blob, data_match.group(1), f"markdown-image-{index}")
            except (ValueError, base64.binascii.Error):
                references.append("无法解码的内嵌图片")
        else:
            references.append(source)
    return text, images, references, {"lineCount": len(text.splitlines()), "referenceImageCount": len(references)}


def _extract_txt(content: bytes) -> Tuple[str, List[Dict[str, Any]], List[str], Dict[str, Any]]:
    text = _decode_text(content)
    return text, [], [], {"lineCount": len(text.splitlines())}


def _extract_docx(content: bytes) -> Tuple[str, List[Dict[str, Any]], List[str], Dict[str, Any]]:
    try:
        from docx import Document
    except ImportError as exc:
        raise FileContentExtractionError("Word 提取依赖未安装") from exc
    try:
        document = Document(io.BytesIO(content))
    except Exception as exc:
        raise FileContentExtractionError(f"Word 文件读取失败: {exc}") from exc

    blocks: List[str] = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                blocks.append(row_text)

    images: List[Dict[str, Any]] = []
    seen_parts = set()
    for relation in document.part.rels.values():
        if "image" not in relation.reltype:
            continue
        part = relation.target_part
        if part.partname in seen_parts:
            continue
        seen_parts.add(part.partname)
        _append_image(images, part.blob, part.content_type, Path(str(part.partname)).name)
    return "\n".join(blocks), images, [], {"paragraphCount": len(document.paragraphs), "tableCount": len(document.tables)}


def _extract_pptx(content: bytes) -> Tuple[str, List[Dict[str, Any]], List[str], Dict[str, Any]]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise FileContentExtractionError("PPT 提取依赖未安装") from exc
    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise FileContentExtractionError(f"PPT 文件读取失败: {exc}") from exc

    blocks: List[str] = []
    images: List[Dict[str, Any]] = []

    def visit_shapes(shapes: Iterable[Any], slide_number: int) -> None:
        for shape in shapes:
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                visit_shapes(shape.shapes, slide_number)
                continue
            if getattr(shape, "has_text_frame", False):
                value = str(shape.text or "").strip()
                if value:
                    blocks.append(value)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        blocks.append(row_text)
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                _append_image(images, image.blob, image.content_type, f"slide-{slide_number}-{image.filename}", slide=slide_number)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        before = len(blocks)
        visit_shapes(slide.shapes, slide_number)
        if len(blocks) > before:
            blocks.insert(before, f"【第 {slide_number} 张幻灯片】")
    return "\n".join(blocks), images, [], {"slideCount": len(presentation.slides)}


def _extract_pdf(content: bytes) -> Tuple[str, List[Dict[str, Any]], List[str], Dict[str, Any]]:
    try:
        import fitz
    except ImportError as exc:
        raise FileContentExtractionError("PDF 提取依赖未安装") from exc
    try:
        document = fitz.open(stream=content, filetype="pdf")
    except Exception as exc:
        raise FileContentExtractionError(f"PDF 文件读取失败: {exc}") from exc

    page_count = document.page_count
    blocks: List[str] = []
    images: List[Dict[str, Any]] = []
    seen_xrefs = set()
    scanned_pages = 0
    try:
        for page_index in range(page_count):
            page = document.load_page(page_index)
            page_number = page_index + 1
            page_text = str(page.get_text("text") or "").strip()
            if page_text:
                blocks.extend((f"【第 {page_number} 页】", page_text))
            page_image_count_before = len(images)
            for image_info in page.get_images(full=True):
                xref = image_info[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                extracted = document.extract_image(xref)
                blob = extracted.get("image", b"")
                extension = str(extracted.get("ext") or "png")
                mime = mimetypes.guess_type(f"image.{extension}")[0] or f"image/{extension}"
                _append_image(images, blob, mime, f"page-{page_number}-image-{xref}.{extension}", page=page_number)
            if not page_text:
                scanned_pages += 1
                if len(images) == page_image_count_before:
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                    _append_image(images, pixmap.tobytes("png"), "image/png", f"page-{page_number}.png", page=page_number)
    finally:
        document.close()
    return "\n".join(blocks), images, [], {"pageCount": page_count, "scannedPageCount": scanned_pages}


def _append_image(images: List[Dict[str, Any]], blob: bytes, mime_type: str, name: str, **location: Any) -> None:
    if not blob or len(images) >= MAX_IMAGES or len(blob) > MAX_IMAGE_BYTES:
        return
    normalized_mime = mime_type if str(mime_type).startswith("image/") else "image/png"
    images.append({
        "name": name,
        "mimeType": normalized_mime,
        "size": len(blob),
        "dataUrl": f"data:{normalized_mime};base64,{base64.b64encode(blob).decode('ascii')}",
        **location,
    })


def _normalize_text(value: str) -> str:
    text = str(value or "").replace("\x00", " ")
    text = re.sub(r"[ \t\r]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()
