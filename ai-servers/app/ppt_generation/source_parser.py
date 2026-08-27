from __future__ import annotations

import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree

from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation


class PptSourceParseError(RuntimeError):
    pass


def extract_source_text(path: Path, max_characters: int = 200_000) -> str:
    extension = path.suffix.lower()
    if extension == ".txt":
        text = _decode_text(path.read_bytes())
    elif extension == ".docx":
        text = _extract_docx(path)
    elif extension == ".pptx":
        text = _extract_pptx(path)
    elif extension == ".xlsx":
        text = _extract_xlsx(path)
    elif extension in {".doc", ".ppt", ".xls"}:
        text = _extract_legacy_office(path)
    else:
        raise PptSourceParseError("不支持的资料格式")
    normalized = re.sub(r"\n{3,}", "\n\n", str(text or "")).strip()
    if not normalized:
        raise PptSourceParseError("未从资料中解析到可用文本")
    return normalized[:max_characters]


def _extract_docx(path: Path) -> str:
    """Extract a DOCX in document order, including headings and tables.

    ``Document.paragraphs`` silently omits tables. PPT source material often
    puts the facts that must become comparison slides in tables, so walking the
    document body is required instead of collecting paragraphs alone.
    """
    document = Document(path)
    blocks = []
    table_index = 0
    for child in document.element.body.iterchildren():
        tag = str(child.tag).rsplit("}", 1)[-1]
        if tag == "p":
            paragraph = Paragraph(child, document)
            value = str(paragraph.text or "").strip()
            if not value:
                continue
            style_name = str(getattr(paragraph.style, "name", "") or "")
            heading = re.search(r"(?:Heading|标题)\s*([1-6])", style_name, flags=re.IGNORECASE)
            if heading:
                value = f"{'#' * int(heading.group(1))} {value}"
            blocks.append(value)
        elif tag == "tbl":
            table_index += 1
            table = Table(child, document)
            rows = []
            for row in table.rows:
                cells = [re.sub(r"\s+", " ", str(cell.text or "").strip()) for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                blocks.append(f"[表格 {table_index}]\n" + "\n".join(rows))
    return "\n".join(blocks)


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _extract_pptx(path: Path) -> str:
    rows = []
    for index, slide in enumerate(Presentation(path).slides, start=1):
        values = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = str(shape.text or "").strip()
                if value:
                    values.append(value)
        if values:
            rows.append(f"第 {index} 页\n" + "\n".join(values))
    return "\n\n".join(rows)


def _extract_xlsx(path: Path) -> str:
    namespace = {"m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    with zipfile.ZipFile(path) as archive:
        shared = []
        if "xl/sharedStrings.xml" in archive.namelist():
            root = ElementTree.fromstring(archive.read("xl/sharedStrings.xml"))
            for item in root.findall("m:si", namespace):
                shared.append("".join(node.text or "" for node in item.findall(".//m:t", namespace)))
        rows = []
        sheets = sorted(name for name in archive.namelist() if name.startswith("xl/worksheets/sheet") and name.endswith(".xml"))
        for sheet_index, name in enumerate(sheets, start=1):
            root = ElementTree.fromstring(archive.read(name))
            rows.append(f"工作表 {sheet_index}")
            for row in root.findall(".//m:row", namespace):
                values = []
                for cell in row.findall("m:c", namespace):
                    value = cell.find("m:v", namespace)
                    raw = value.text if value is not None else ""
                    if cell.get("t") == "s" and raw.isdigit() and int(raw) < len(shared):
                        raw = shared[int(raw)]
                    if raw:
                        values.append(raw)
                if values:
                    rows.append(" | ".join(values))
        return "\n".join(rows)


def _extract_legacy_office(path: Path) -> str:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise PptSourceParseError("旧版 Office 文件解析需要安装 LibreOffice")
    target_extension = "pptx" if path.suffix.lower() == ".ppt" else "xlsx" if path.suffix.lower() == ".xls" else "docx"
    with tempfile.TemporaryDirectory(prefix="a3-ppt-source-") as directory:
        subprocess.run(
            [soffice, "--headless", "--convert-to", target_extension, "--outdir", directory, str(path)],
            check=True,
            capture_output=True,
            timeout=120,
        )
        converted = Path(directory) / f"{path.stem}.{target_extension}"
        if not converted.is_file():
            raise PptSourceParseError("LibreOffice 未生成可解析文件")
        return extract_source_text(converted)

