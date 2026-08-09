from __future__ import annotations

import copy
import json
import math
import re
from collections import deque
from datetime import date, datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Mapping, MutableMapping, Sequence, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.ppt_generation.template_catalog import EmbeddedTemplateCatalog
from app.rag.document_conversion import generated_exporter


_CANVAS_WIDTH = 1280.0
_CANVAS_HEIGHT = 720.0
_SLIDE_WIDTH = 13.333
_SLIDE_HEIGHT = 7.5
_X_SCALE = _SLIDE_WIDTH / _CANVAS_WIDTH
_Y_SCALE = _SLIDE_HEIGHT / _CANVAS_HEIGHT
_TEMPLATE_ROOT = Path(__file__).resolve().parent / "assets" / "templates"
_REPLACEABLE_IMAGE = "replaceable_template_image"

_TITLE_WORDS = ("title", "heading", "headline", "section_title", "cover_title", "main_title")
_BODY_WORDS = (
    "body", "copy", "paragraph", "description", "caption", "subtitle", "tagline",
    "detail", "supporting", "summary", "quote", "insight",
)
_PAGE_WORDS = ("page", "pagination", "folio", "footer_index", "slide_marker")
_LABEL_WORDS = ("label", "item_title", "card_title", "feature_title", "metric_title")


def render_presenton_presentation(
    slides: Iterable[Any],
    title: str,
    settings: Mapping[str, Any],
) -> Tuple[Dict[str, Any], Path]:
    """Render an editable PPTX directly from a bundled Presenton template JSON."""
    specs = [_normalize_slide(item, index) for index, item in enumerate(slides, start=1)]
    if len(specs) < 2:
        raise ValueError("PPT 至少需要两页")

    template_id = str(settings.get("templateId") or "general").strip().lower()
    catalog = EmbeddedTemplateCatalog()
    if not catalog.contains(template_id):
        raise ValueError(f"PPT 模板不存在：{template_id}")
    template = catalog.load(template_id)
    layouts = [item for item in template.get("layouts") or [] if isinstance(item, dict)]
    if not layouts:
        raise ValueError(f"PPT 模板没有可用布局：{template_id}")

    presentation = Presentation()
    presentation.slide_width = Inches(_SLIDE_WIDTH)
    presentation.slide_height = Inches(_SLIDE_HEIGHT)
    presentation.core_properties.title = title
    presentation.core_properties.subject = f"AgentA3 内置 Presenton 模板：{template_id}"
    presentation.core_properties.created = datetime(2000, 1, 1)
    presentation.core_properties.modified = datetime(2000, 1, 1)

    selected_layouts: List[str] = []
    for index, spec in enumerate(specs, start=1):
        layout = _select_layout(layouts, spec, index, len(specs), selected_layouts)
        selected_layouts.append(str(layout.get("id") or ""))
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        binder = _ContentBinder(spec, title, index, len(specs), _template_accent(template))
        for component in layout.get("components") or []:
            _render_node(
                slide,
                copy.deepcopy(component),
                binder,
                template_id,
                origin=(0.0, 0.0),
                inherited_size=(_CANVAS_WIDTH, _CANVAS_HEIGHT),
            )

    path = generated_exporter._new_export_path(generated_exporter._slugify(title), "pptx")
    generated_exporter._atomic_write_payload(path, presentation.save)
    attachment = generated_exporter._attachment_for_file(
        path,
        "ai_ppt_generation_tool",
        "PPTX",
        display_stem=title,
    )
    attachment.update({
        "type": "pptx",
        "templateId": template_id,
        "templateLayoutIds": selected_layouts,
    })
    return attachment, path


def _normalize_slide(value: Any, index: int) -> Dict[str, Any]:
    item = value if isinstance(value, Mapping) else {}
    content = item.get("content") or []
    if isinstance(content, str):
        points = [line.strip(" -*•") for line in content.splitlines() if line.strip(" -*•")]
    elif isinstance(content, list):
        points = [str(point).strip() for point in content if str(point).strip()]
    else:
        points = []
    return {
        "index": index,
        "title": str(item.get("title") or f"第 {index} 页").strip()[:160],
        "content": points[:12],
        "objective": str(item.get("objective") or "").strip(),
        "type": str(item.get("type") or ("cover" if index == 1 else "content")).strip(),
        "layout": str(item.get("templateLayoutId") or item.get("layout") or "").strip(),
        "visualPrompt": str(item.get("visualPrompt") or "").strip(),
        "imagePath": str(item.get("imagePath") or "").strip(),
    }


def _select_layout(
    layouts: Sequence[Dict[str, Any]],
    spec: Mapping[str, Any],
    index: int,
    total: int,
    previous: Sequence[str],
) -> Dict[str, Any]:
    requested = str(spec.get("layout") or "").strip().lower()
    for layout in layouts:
        if requested and str(layout.get("id") or "").lower() == requested:
            return layout

    slide_type = str(spec.get("type") or "").lower()
    text = " ".join([slide_type, str(spec.get("title") or ""), str(spec.get("visualPrompt") or "")]).lower()
    content = spec.get("content") if isinstance(spec.get("content"), list) else []
    has_numbers = any(re.search(r"\d", str(value)) for value in content)
    has_table = any("|" in str(value) or "表" in str(value) for value in content)
    has_image = bool(spec.get("imagePath"))

    wanted: List[str]
    if index == 1 or any(word in text for word in ("cover", "封面", "开场")):
        wanted = ["cover", "title_intro", "title_description", "hero"]
    elif any(word in text for word in ("目录", "contents", "agenda")):
        wanted = ["table_of_contents", "contents", "agenda"]
    elif any(word in text for word in ("章节", "过渡", "section")):
        wanted = ["section", "quote", "title_description"]
    elif has_table or any(word in text for word in ("对比", "比较", "table")):
        wanted = ["table", "comparison", "compare"]
    elif has_numbers or any(word in text for word in ("数据", "指标", "趋势", "chart", "metric")):
        wanted = ["chart", "metric", "data"]
    elif any(word in text for word in ("流程", "步骤", "过程", "时间线", "process", "timeline")):
        wanted = ["process", "timeline", "steps", "bullet_points"]
    elif any(word in text for word in ("总结", "结论", "回顾", "summary", "conclusion")):
        wanted = ["summary", "quote", "metrics", "description"]
    elif len(content) >= 5:
        wanted = ["bullet_points", "cards", "grid", "metrics_description"]
    elif has_image:
        wanted = ["image_description", "image_bullet", "image"]
    else:
        wanted = ["description_bullet", "bullet_points", "description", "cards"]

    best = layouts[0]
    best_score = -10_000
    for position, layout in enumerate(layouts):
        layout_id = str(layout.get("id") or "").lower()
        description = str(layout.get("description") or "").lower()
        haystack = f"{layout_id} {description}"
        score = sum(8 - min(rank, 6) for rank, word in enumerate(wanted) if word in haystack)
        kinds = _element_types(layout)
        if "image" in kinds and not has_image:
            score -= 3
        if _contains_replaceable_image(layout) and not has_image:
            score -= 1
        if "chart" in kinds and not has_numbers:
            score -= 9
        if "table" in kinds and not has_table:
            score -= 8
        if previous and layout_id == previous[-1]:
            score -= 2
        score -= position * 0.001
        if score > best_score:
            best, best_score = layout, score
    return best


def _element_types(value: Any) -> set[str]:
    result: set[str] = set()
    if isinstance(value, list):
        for item in value:
            result.update(_element_types(item))
    elif isinstance(value, dict):
        if value.get("type"):
            result.add(str(value["type"]))
        for key in ("components", "elements", "children", "child"):
            result.update(_element_types(value.get(key)))
    return result


def _contains_replaceable_image(value: Any) -> bool:
    if isinstance(value, list):
        return any(_contains_replaceable_image(item) for item in value)
    if not isinstance(value, dict):
        return False
    if value.get("type") == "image" and _REPLACEABLE_IMAGE in str(value.get("data") or ""):
        return True
    return any(_contains_replaceable_image(value.get(key)) for key in ("components", "elements", "children", "child"))


class _ContentBinder:
    def __init__(
        self,
        spec: Mapping[str, Any],
        deck_title: str,
        page: int,
        total: int,
        accent: RGBColor,
    ) -> None:
        self.title = str(spec.get("title") or deck_title)
        values = [str(value).strip() for value in spec.get("content") or [] if str(value).strip()]
        objective = str(spec.get("objective") or "").strip()
        if objective and objective not in values:
            values.append(objective)
        self.values = values or [self.title]
        self.pending = deque(self.values)
        self.page = page
        self.total = total
        self.image_path = str(spec.get("imagePath") or "").strip()
        self.accent = accent

    def text(self, node: Mapping[str, Any]) -> str:
        name = str(node.get("name") or "").lower()
        sample = "".join(str(run.get("text") or "") for run in node.get("runs") or [] if isinstance(run, Mapping))
        if node.get("decorative"):
            return sample
        if any(word in name for word in _PAGE_WORDS):
            return f"{self.page:02d}" if "total" not in name else f"{self.page:02d} / {self.total:02d}"
        if "date" in name:
            return date.today().isoformat()
        if "presenter" in name or "author" in name or "attribution" in name or "avatar" in name:
            return ""
        if any(word in name for word in _TITLE_WORDS) and not any(word in name for word in ("item", "card", "feature")):
            return self.title
        if any(word in name for word in _BODY_WORDS + _LABEL_WORDS):
            return self._next()
        return self._next()

    def list_items(self, maximum: int = 8) -> List[str]:
        return self.values[:maximum]

    def _next(self) -> str:
        return self.pending.popleft() if self.pending else ""


def _render_node(
    slide,
    node: Any,
    binder: _ContentBinder,
    template_id: str,
    *,
    origin: Tuple[float, float],
    inherited_size: Tuple[float, float],
) -> None:
    if not isinstance(node, MutableMapping):
        return
    position = node.get("position") if isinstance(node.get("position"), Mapping) else {}
    x = origin[0] + float(position.get("x") or 0)
    y = origin[1] + float(position.get("y") or 0)
    size = node.get("size") if isinstance(node.get("size"), Mapping) else {}
    width = float(size.get("width") or inherited_size[0])
    height = float(size.get("height") or inherited_size[1])
    kind = str(node.get("type") or "component")

    if kind in {"component", "group"}:
        for child in node.get("elements") or node.get("children") or []:
            _render_node(slide, child, binder, template_id, origin=(x, y), inherited_size=(width, height))
        return
    if kind == "container":
        child = node.get("child")
        if child:
            _render_node(slide, child, binder, template_id, origin=(x, y), inherited_size=(width, height))
        return
    if kind in {"flex", "grid"}:
        _render_flow(slide, node, binder, template_id, (x, y), (width, height), kind)
        return
    if kind == "text":
        _render_text(slide, node, binder.text(node), x, y, width, height)
    elif kind == "text-list":
        _render_text_list(slide, node, binder.list_items(int(node.get("max_items") or 8)), x, y, width, height)
    elif kind == "image":
        _render_image(slide, node, binder, template_id, x, y, width, height)
    elif kind == "vector":
        _render_vector(slide, node, x, y, width, height)
    elif kind == "chart":
        _render_chart(slide, node, binder, x, y, width, height)
    elif kind == "table":
        _render_table(slide, node, binder, x, y, width, height)


def _render_flow(slide, node, binder, template_id, origin, size, kind: str) -> None:
    children = [copy.deepcopy(item) for item in node.get("children") or [] if isinstance(item, dict)]
    if not children:
        return
    maximum = int(node.get("max_children") or len(children))
    minimum = int(node.get("min_children") or len(children))
    wanted = max(minimum, min(maximum, len(binder.values)))
    if len(children) == 1 and wanted > 1:
        children = [copy.deepcopy(children[0]) for _ in range(wanted)]
    gap = float(node.get("gap") or 0)
    if kind == "grid":
        columns = max(1, int(node.get("columns") or math.ceil(math.sqrt(len(children)))))
        rows = max(1, int(node.get("rows") or math.ceil(len(children) / columns)))
        cell_w = max(1.0, (size[0] - gap * (columns - 1)) / columns)
        cell_h = max(1.0, (size[1] - gap * (rows - 1)) / rows)
        for index, child in enumerate(children):
            row, column = divmod(index, columns)
            child["position"] = {"x": column * (cell_w + gap), "y": row * (cell_h + gap)}
            child["size"] = {"width": cell_w, "height": cell_h}
            _render_node(slide, child, binder, template_id, origin=origin, inherited_size=(cell_w, cell_h))
        return
    direction = str(node.get("direction") or "row")
    cursor = 0.0
    for child in children:
        child_size = child.get("size") if isinstance(child.get("size"), Mapping) else {}
        child_w = float(child_size.get("width") or (size[0] / max(1, len(children))))
        child_h = float(child_size.get("height") or (size[1] / max(1, len(children))))
        child["position"] = {"x": cursor if direction == "row" else 0, "y": cursor if direction != "row" else 0}
        _render_node(slide, child, binder, template_id, origin=origin, inherited_size=(child_w, child_h))
        cursor += (child_w if direction == "row" else child_h) + gap


def _render_text(slide, node, value: str, x: float, y: float, width: float, height: float) -> None:
    if not value:
        return
    style = _text_style(node)
    minimum_height = (style["size"] / 0.75) * 1.35
    if y < 80 and style["size"] >= 30:
        y += (style["size"] / 0.75) * 0.22
    box = slide.shapes.add_textbox(
        _x(x), _y(y), _x(max(width, 4)), _y(max(height, minimum_height, 4)),
    )
    box.rotation = float(node.get("rotation") or 0)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    alignment = node.get("alignment") if isinstance(node.get("alignment"), Mapping) else {}
    frame.vertical_anchor = {
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(str(alignment.get("vertical") or "top"), MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(str(alignment.get("horizontal") or "left"), PP_ALIGN.LEFT)
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = _fit_text(value, node)
    run.font.name = "Noto Sans CJK SC" if re.search(r"[\u3400-\u9fff]", value) else style["family"]
    run.font.size = Pt(style["size"])
    run.font.bold = style["bold"]
    run.font.italic = style["italic"]
    run.font.color.rgb = style["color"]


def _render_text_list(slide, node, items, x, y, width, height) -> None:
    if not items:
        return
    box = slide.shapes.add_textbox(_x(x), _y(y), _x(width), _y(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Pt(2)
    frame.margin_top = frame.margin_bottom = 0
    style = _text_style(node)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Noto Sans CJK SC" if re.search(r"[\u3400-\u9fff]", item) else style["family"]
        paragraph.font.size = Pt(style["size"])
        paragraph.font.color.rgb = style["color"]
        paragraph.space_after = Pt(5)
        paragraph.text = f"• {item}"


def _render_image(slide, node, binder, template_id, x, y, width, height) -> None:
    raw = str(node.get("data") or "").strip()
    path: Path | None = None
    if _REPLACEABLE_IMAGE in raw:
        if binder.image_path:
            candidate = Path(binder.image_path).resolve()
            if candidate.is_file():
                path = candidate
    else:
        normalized = raw.lstrip("/")
        if normalized.startswith("static/images/"):
            normalized = normalized[len("static/images/"):]
        elif normalized.startswith("static/"):
            normalized = normalized[len("static/"):]
        candidate = (_TEMPLATE_ROOT / template_id / "static" / normalized).resolve()
        template_root = (_TEMPLATE_ROOT / template_id / "static").resolve()
        if candidate.is_relative_to(template_root) and candidate.is_file() and candidate.suffix.lower() != ".svg":
            path = candidate
    if path is None:
        if _REPLACEABLE_IMAGE in raw:
            placeholder = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                _x(x), _y(y), _x(max(width, 1)), _y(max(height, 1)),
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = binder.accent
            placeholder.fill.transparency = 72
            placeholder.line.fill.background()
        return
    slide.shapes.add_picture(str(path), _x(x), _y(y), width=_x(width), height=_y(height))


def _render_vector(slide, node, x, y, width, height) -> None:
    points = node.get("points") if isinstance(node.get("points"), list) else []
    shape_name = str(node.get("shape") or "polygon")
    if points:
        xs = [float(item.get("x") or 0) for item in points if isinstance(item, Mapping)]
        ys = [float(item.get("y") or 0) for item in points if isinstance(item, Mapping)]
        if not xs or not ys:
            return
        if not node.get("closed") and len(xs) == 2:
            shape = slide.shapes.add_connector(1, _x(x + xs[0]), _y(y + ys[0]), _x(x + xs[1]), _y(y + ys[1]))
            _apply_line(shape, node.get("stroke"))
            return
        if shape_name == "ellipse":
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, _x(x + min(xs)), _y(y + min(ys)), _x(max(xs) - min(xs)), _y(max(ys) - min(ys)))
        else:
            builder = slide.shapes.build_freeform(xs[0], ys[0], scale=(_x(1), _y(1)))
            builder.add_line_segments(list(zip(xs[1:], ys[1:])), close=bool(node.get("closed", True)))
            shape = builder.convert_to_shape(_x(x), _y(y))
    else:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL if shape_name == "ellipse" else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            _x(x), _y(y), _x(max(width, 1)), _y(max(height, 1)),
        )
    _apply_fill(shape, node.get("fill"))
    _apply_line(shape, node.get("stroke"))
    if node.get("rotation"):
        shape.rotation = float(node["rotation"])


def _render_chart(slide, node, binder, x, y, width, height) -> None:
    values: List[float] = []
    labels: List[str] = []
    for item in binder.values:
        match = re.search(r"(-?\d+(?:\.\d+)?)", item.replace(",", ""))
        if match:
            labels.append(re.sub(r"[-+]?\d+(?:\.\d+)?%?", "", item).strip(" ：:-") or f"项目 {len(labels) + 1}")
            values.append(float(match.group(1)))
    if len(values) < 2:
        return
    data = CategoryChartData()
    data.categories = labels[:8]
    data.add_series("数据", values[:8])
    chart_type = {
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    }.get(str(node.get("chart_type") or "column").lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    chart = slide.shapes.add_chart(chart_type, _x(x), _y(y), _x(width), _y(height), data).chart
    chart.has_legend = bool(node.get("legend"))
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.has_title = False


def _render_table(slide, node, binder, x, y, width, height) -> None:
    rows = binder.list_items(int(node.get("max_rows") or 6))
    if not rows:
        return
    parsed = [[part.strip() for part in value.strip("|").split("|")] for value in rows]
    column_count = max(2, max(len(row) for row in parsed))
    if column_count == 2 and all(len(row) == 1 for row in parsed):
        parsed = [[str(index + 1), row[0]] for index, row in enumerate(parsed)]
    table = slide.shapes.add_table(len(parsed), column_count, _x(x), _y(y), _x(width), _y(height)).table
    font = _text_style(node)
    for row_index, values in enumerate(parsed):
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            cell.text = values[column_index] if column_index < len(values) else ""
            cell.margin_left = cell.margin_right = Pt(4)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = font["family"]
                paragraph.font.size = Pt(font["size"])
                paragraph.font.color.rgb = font["color"]


def _text_style(node: Mapping[str, Any]) -> Dict[str, Any]:
    font = node.get("font") if isinstance(node.get("font"), Mapping) else {}
    runs = node.get("runs") if isinstance(node.get("runs"), list) else []
    if runs and isinstance(runs[0], Mapping) and isinstance(runs[0].get("font"), Mapping):
        font = {**font, **runs[0]["font"]}
    family = str(font.get("family") or "Noto Sans CJK SC").replace(" Semi-Bold", "").strip()
    return {
        "family": family,
        "size": max(8.0, float(font.get("size") or 18) * 0.75),
        "bold": bool(font.get("bold")),
        "italic": bool(font.get("italic")),
        "color": _color(font.get("color"), "172033"),
    }


def _fit_text(value: str, node: Mapping[str, Any]) -> str:
    maximum = int(node.get("max_length") or 0)
    if maximum > 0 and len(value) > maximum:
        return value[: max(1, maximum - 1)].rstrip("，。；;、 ") + "…"
    return value


def _apply_fill(shape, fill: Any) -> None:
    if not isinstance(fill, Mapping) or not fill.get("color"):
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = _color(fill.get("color"), "FFFFFF")
    opacity = float(fill.get("opacity", 1) or 0)
    shape.fill.transparency = int(max(0, min(100, (1 - opacity) * 100)))


def _apply_line(shape, stroke: Any) -> None:
    if not isinstance(stroke, Mapping) or not stroke.get("color"):
        shape.line.fill.background()
        return
    shape.line.color.rgb = _color(stroke.get("color"), "000000")
    shape.line.width = Pt(max(0.25, float(stroke.get("width") or 1) * 0.75))


def _color(value: Any, fallback: str) -> RGBColor:
    raw = re.sub(r"[^0-9A-Fa-f]", "", str(value or fallback))
    if len(raw) == 3:
        raw = "".join(char * 2 for char in raw)
    if len(raw) == 8:
        raw = raw[-6:]
    if len(raw) != 6:
        raw = fallback
    return RGBColor.from_string(raw.upper())


def _template_accent(template: Mapping[str, Any]) -> RGBColor:
    counts: Dict[str, int] = {}

    def visit(value: Any) -> None:
        if isinstance(value, list):
            for item in value:
                visit(item)
            return
        if not isinstance(value, Mapping):
            return
        for key in ("fill", "stroke", "font"):
            style = value.get(key)
            if isinstance(style, Mapping) and style.get("color"):
                raw = re.sub(r"[^0-9A-Fa-f]", "", str(style["color"])).upper()
                if len(raw) == 6 and raw not in {"FFFFFF", "000000", "111111", "111827", "172033", "F8FAFC"}:
                    counts[raw] = counts.get(raw, 0) + 1
        for key in ("layouts", "components", "elements", "children", "child"):
            visit(value.get(key))

    visit(template.get("layouts") or [])
    def accent_score(raw: str) -> float:
        red, green, blue = (int(raw[index:index + 2], 16) for index in (0, 2, 4))
        maximum = max(red, green, blue)
        minimum = min(red, green, blue)
        saturation = 0.0 if maximum == 0 else (maximum - minimum) / maximum
        return saturation * 100 + min(counts[raw], 20)

    chosen = max(counts, key=accent_score) if counts else "5265F5"
    return RGBColor.from_string(chosen)


def _x(value: float):
    return Inches(max(0.0, float(value)) * _X_SCALE)


def _y(value: float):
    return Inches(max(0.0, float(value)) * _Y_SCALE)


__all__ = ["render_presenton_presentation"]
