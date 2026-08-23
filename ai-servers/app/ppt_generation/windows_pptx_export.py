"""Small Windows-native PPTX exporter for the embedded Presenton UI tree.

The official Presenton converter ships a Linux executable. HBuilder users run
the AI server directly on Windows, so this exporter keeps the generated deck
usable there without introducing a non-editable export fallback. Text boxes remain
editable and the template's positions, fonts, colours and basic vector panels
are carried over from the already validated UI tree.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Mapping

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SLIDE_WIDTH = 1280.0
SLIDE_HEIGHT = 720.0
PPT_WIDTH = 13.333
PPT_HEIGHT = 7.5


def export_presenton_slides(slides: list[Mapping[str, Any]], path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(PPT_WIDTH)
    presentation.slide_height = Inches(PPT_HEIGHT)
    blank = presentation.slide_layouts[6]
    for spec in slides:
        slide = presentation.slides.add_slide(blank)
        ui = spec.get("ui") if isinstance(spec, Mapping) else None
        if isinstance(ui, Mapping):
            background = ui.get("background")
            _set_background(slide, background)
            for key in ("components", "elements"):
                _walk(slide, ui.get(key), 0.0, 0.0)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _walk(slide: Any, value: Any, parent_x: float, parent_y: float) -> None:
    if isinstance(value, list):
        for item in value:
            _walk(slide, item, parent_x, parent_y)
        return
    if not isinstance(value, Mapping):
        return
    position = value.get("position") if isinstance(value.get("position"), Mapping) else {}
    offset_x = parent_x + _number(position.get("x"))
    offset_y = parent_y + _number(position.get("y"))
    node_type = str(value.get("type") or "").lower()
    size = value.get("size") if isinstance(value.get("size"), Mapping) else {}
    width = _number(size.get("width"))
    height = _number(size.get("height"))
    if node_type in {"vector", "shape", "rect", "circle", "line", "divider"}:
        _add_vector(slide, value, offset_x, offset_y, width, height)
    elif node_type == "text" and width > 0 and height > 0:
        _add_text(slide, value, offset_x, offset_y, width, height)
    elif node_type == "image" and width > 0 and height > 0:
        _add_image(slide, value, offset_x, offset_y, width, height)
    if node_type == "flex":
        _walk_flex_children(slide, value, offset_x, offset_y, width, height)
        child_keys = ("components", "elements")
    elif node_type == "grid":
        _walk_grid_children(slide, value, offset_x, offset_y, width, height)
        child_keys = ("components", "elements")
    else:
        child_keys = ("components", "elements", "children")
    for key in child_keys:
        _walk(slide, value.get(key), offset_x, offset_y)
    _walk(slide, value.get("child"), offset_x, offset_y)


def _walk_flex_children(
    slide: Any,
    node: Mapping[str, Any],
    parent_x: float,
    parent_y: float,
    width: float,
    height: float,
) -> None:
    """Apply the small flex subset used by Presenton template UI trees.

    The browser renderer lays out children such as stacked cards, while a
    naive PPTX tree walk would place every child at the flex container's
    origin. Keep the exporter deterministic and editable by resolving the
    common row/column direction, gap, alignment and justification properties.
    """
    children = node.get("children")
    if not isinstance(children, list) or not children:
        return
    direction = str(node.get("direction") or "row").lower()
    is_column = direction == "column"
    gap = max(0.0, _number(node.get("gap")))
    sizes = [
        child.get("size") if isinstance(child, Mapping) and isinstance(child.get("size"), Mapping) else {}
        for child in children
    ]
    main_sizes = [
        _number(size.get("height" if is_column else "width"))
        for size in sizes
    ]
    cross_sizes = [
        _number(size.get("width" if is_column else "height"))
        for size in sizes
    ]
    main_extent = height if is_column else width
    cross_extent = width if is_column else height
    occupied = sum(main_sizes) + gap * max(0, len(children) - 1)
    justify = str(node.get("justify_content") or "flex-start").lower().replace("-", "_")
    if justify in {"center", "middle"}:
        cursor = max(0.0, (main_extent - occupied) / 2.0)
        actual_gap = gap
    elif justify in {"flex_end", "end"}:
        cursor = max(0.0, main_extent - occupied)
        actual_gap = gap
    elif justify == "space_between" and len(children) > 1:
        cursor = 0.0
        actual_gap = max(gap, (main_extent - sum(main_sizes)) / (len(children) - 1))
    else:
        cursor = 0.0
        actual_gap = gap
    align = str(node.get("align_items") or "flex-start").lower().replace("-", "_")
    for child, size, main_size, cross_size in zip(children, sizes, main_sizes, cross_sizes):
        if not isinstance(child, Mapping):
            cursor += main_size + actual_gap
            continue
        if align in {"center", "middle"}:
            cross_cursor = max(0.0, (cross_extent - cross_size) / 2.0)
        elif align in {"flex_end", "end"}:
            cross_cursor = max(0.0, cross_extent - cross_size)
        else:
            cross_cursor = 0.0
        if is_column:
            layout_x, layout_y = cross_cursor, cursor
        else:
            layout_x, layout_y = cursor, cross_cursor
        explicit = child.get("position") if isinstance(child.get("position"), Mapping) else {}
        explicit_x = _number(explicit.get("x"))
        explicit_y = _number(explicit.get("y"))
        _walk(
            slide,
            child,
            parent_x + layout_x - explicit_x,
            parent_y + layout_y - explicit_y,
        )
        cursor += main_size + actual_gap


def _walk_grid_children(
    slide: Any,
    node: Mapping[str, Any],
    parent_x: float,
    parent_y: float,
    width: float,
    height: float,
) -> None:
    """Resolve the equal-column grid used by metric/card template blocks."""
    children = node.get("children")
    if not isinstance(children, list) or not children:
        return
    try:
        columns = max(1, int(node.get("columns") or 1))
    except (TypeError, ValueError):
        columns = 1
    try:
        rows = max(1, int(node.get("rows") or ((len(children) + columns - 1) // columns)))
    except (TypeError, ValueError):
        rows = max(1, (len(children) + columns - 1) // columns)
    column_gap = max(0.0, _number(node.get("column_gap")))
    row_gap = max(0.0, _number(node.get("row_gap")))
    cell_width = max(0.0, (width - column_gap * max(0, columns - 1)) / columns)
    cell_height = max(0.0, (height - row_gap * max(0, rows - 1)) / rows)
    for index, child in enumerate(children):
        if not isinstance(child, Mapping):
            continue
        row, column = divmod(index, columns)
        layout_x = column * (cell_width + column_gap)
        layout_y = row * (cell_height + row_gap)
        explicit = child.get("position") if isinstance(child.get("position"), Mapping) else {}
        explicit_x = _number(explicit.get("x"))
        explicit_y = _number(explicit.get("y"))
        _walk(
            slide,
            child,
            parent_x + layout_x - explicit_x,
            parent_y + layout_y - explicit_y,
        )


def _add_text(slide: Any, node: Mapping[str, Any], x: float, y: float, width: float, height: float) -> None:
    shape = slide.shapes.add_textbox(_inches(x), _inches(y), _inches(width), _inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = {
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(str((node.get("alignment") or {}).get("vertical") or "top").lower(), MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    horizontal = str((node.get("alignment") or {}).get("horizontal") or "left").lower()
    paragraph.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(horizontal, PP_ALIGN.LEFT)
    text = str(node.get("text") or "")
    if not text:
        runs = node.get("runs")
        if isinstance(runs, list):
            text = "".join(str(run.get("text") or "") for run in runs if isinstance(run, Mapping))
    run = paragraph.add_run()
    run.text = text
    font = node.get("font") if isinstance(node.get("font"), Mapping) else {}
    run.font.name = str(font.get("family") or "Arial")
    run.font.size = Pt(max(7.0, _number(font.get("size")) * 0.75))
    run.font.bold = bool(font.get("bold"))
    run.font.italic = bool(font.get("italic"))
    run.font.color.rgb = _rgb(font.get("color"), RGBColor(17, 24, 39))


def _add_vector(slide: Any, node: Mapping[str, Any], x: float, y: float, width: float, height: float) -> None:
    if width <= 0 or height <= 0:
        points = node.get("points") if isinstance(node.get("points"), list) else []
        if points:
            xs = [_number(point.get("x")) for point in points if isinstance(point, Mapping)]
            ys = [_number(point.get("y")) for point in points if isinstance(point, Mapping)]
            if xs and ys:
                width, height = max(xs) - min(xs), max(ys) - min(ys)
    if width <= 0 or height <= 0:
        return
    shape_name = str(node.get("shape") or node.get("type") or "").lower()
    shape_type = MSO_AUTO_SHAPE_TYPE.OVAL if "circle" in shape_name or "ellipse" in shape_name else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, _inches(x), _inches(y), _inches(width), _inches(height))
    fill = node.get("fill") if isinstance(node.get("fill"), Mapping) else {}
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill.get("color"), RGBColor(255, 255, 255))
    shape.line.fill.background()


def _add_image(slide: Any, node: Mapping[str, Any], x: float, y: float, width: float, height: float) -> None:
    data = str(node.get("data") or "")
    if data.startswith("file://"):
        data = data[7:]
    source = Path(data)
    if source.is_file():
        slide.shapes.add_picture(str(source), _inches(x), _inches(y), _inches(width), _inches(height))


def _set_background(slide: Any, value: Any) -> None:
    if not isinstance(value, str) or not value.startswith("#"):
        return
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(value, RGBColor(255, 255, 255))


def _inches(value: float) -> Inches:
    return Inches(max(0.0, value) / SLIDE_WIDTH * PPT_WIDTH)


def _number(value: Any) -> float:
    try:
        return float(value or 0)
    except (TypeError, ValueError):
        return 0.0


def _rgb(value: Any, fallback: RGBColor) -> RGBColor:
    raw = str(value or "").strip().lstrip("#")
    if len(raw) == 6:
        try:
            return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
        except ValueError:
            pass
    return fallback
